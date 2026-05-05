#!/usr/bin/env python3
from __future__ import annotations

import argparse
import json
import logging
import math
import re
from datetime import datetime
from pathlib import Path
import shutil
import subprocess
import tempfile

from PIL import Image, ImageDraw, ImageFilter, ImageFont, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from lxml import etree

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Auto-output directory helpers
# ---------------------------------------------------------------------------

_SANITIZE_RE = re.compile(r'[<>:"/\\|?*\x00-\x1f]')


def sanitize_dirname(text: str, max_length: int = 40) -> str:
    cleaned = text.strip()
    cleaned = _SANITIZE_RE.sub("", cleaned)
    cleaned = re.sub(r"\s+", "_", cleaned)
    cleaned = cleaned.strip("._")
    cleaned = re.sub(r"_{2,}", "_", cleaned)
    if len(cleaned) > max_length:
        cleaned = cleaned[:max_length].rstrip("._")
    return cleaned or "cyberpunk_deck"


def resolve_output_dir(title: str, base_dir: Path | None = None) -> Path:
    base = base_dir or Path.home() / "ai-gen-ppt"
    ts = datetime.now().strftime("%Y%m%d_%H%M")
    safe = sanitize_dirname(title)
    out_dir = base / f"{safe}_{ts}"
    out_dir.mkdir(parents=True, exist_ok=True)
    return out_dir


def extract_deck_title(spec: dict) -> str:
    if spec.get("deck_title"):
        return spec["deck_title"]
    slides = spec.get("slides", [])
    if slides:
        first = slides[0]
        titles = first.get("title", [])
        if titles:
            parts = [t["text"] for t in titles[:3]]
            return "".join(parts)
        if first.get("ghost"):
            return first["ghost"]
    return "Cyberpunk Deck"

NSMAP = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
}


def _ensure_effect_lst(spPr) -> etree._Element:
    effectLst = spPr.find("a:effectLst", NSMAP)
    if effectLst is None:
        effectLst = etree.SubElement(spPr, "{%s}effectLst" % NSMAP["a"])
    return effectLst


def add_glow_to_shape(shape, glow_color: RGBColor, size: int = 40000) -> None:
    try:
        spPr = shape._element.find(".//a:spPr", NSMAP)
        if spPr is None:
            return
        effectLst = _ensure_effect_lst(spPr)
        glow = etree.SubElement(effectLst, "{%s}glow" % NSMAP["a"])
        glow.set("rad", str(size))
        srgb = etree.SubElement(glow, "{%s}srgbClr" % NSMAP["a"])
        srgb.set("val", "%02X%02X%02X" % (glow_color[0], glow_color[1], glow_color[2]))
        alpha = etree.SubElement(srgb, "{%s}alpha" % NSMAP["a"])
        alpha.set("val", "35000")
    except Exception:
        logger.warning("add_glow_to_shape failed", exc_info=True)


def add_glow_to_run(run, glow_color: RGBColor, size: int = 50000) -> None:
    try:
        rPr = run._r.find(".//a:rPr", NSMAP)
        if rPr is None:
            return
        effectLst = _ensure_effect_lst(rPr)
        glow = etree.SubElement(effectLst, "{%s}glow" % NSMAP["a"])
        glow.set("rad", str(size))
        srgb = etree.SubElement(glow, "{%s}srgbClr" % NSMAP["a"])
        srgb.set("val", "%02X%02X%02X" % (glow_color[0], glow_color[1], glow_color[2]))
        alpha = etree.SubElement(srgb, "{%s}alpha" % NSMAP["a"])
        alpha.set("val", "40000")
    except Exception:
        logger.warning("add_glow_to_run failed", exc_info=True)


def add_outer_shadow(shape, color_rgb: str = "000000",
                     blur_rad: int = 76200, dist: int = 25400,
                     direction: int = 5400000, alpha_pct: int = 40000):
    """Add an outer drop shadow to a shape via OOXML injection."""
    try:
        spPr = shape._element.find(".//a:spPr", NSMAP)
        if spPr is None:
            return
        effectLst = _ensure_effect_lst(spPr)
        outerShdw = etree.SubElement(effectLst, "{%s}outerShdw" % NSMAP["a"])
        outerShdw.set("blurRad", str(blur_rad))
        outerShdw.set("dist", str(dist))
        outerShdw.set("dir", str(direction))
        outerShdw.set("algn", "bl")
        outerShdw.set("rotWithShape", "0")
        srgbClr = etree.SubElement(outerShdw, "{%s}srgbClr" % NSMAP["a"])
        srgbClr.set("val", color_rgb)
        alpha = etree.SubElement(srgbClr, "{%s}alpha" % NSMAP["a"])
        alpha.set("val", str(alpha_pct))
    except Exception:
        logger.warning("add_outer_shadow failed", exc_info=True)


def add_accent_line(slide, left_px, top_px, width_px, color_name, thickness=3):
    """Add a thin horizontal accent/separator line with glow."""
    style = active_style()
    accent = color(color_name)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(left_px), px(top_px), px(width_px), px(thickness))
    line.fill.solid()
    line.fill.fore_color.rgb = accent
    line.line.fill.background()
    add_glow_to_shape(line, accent, size=style["line_glow"])


def add_gradient_panel(slide, left_px, top_px, width_px, height_px, accent_name, transparency=0.30):
    """Add a rounded rectangle card with gradient fill and glow border."""
    style = active_style()
    accent = color(accent_name)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px(left_px), px(top_px), px(width_px), px(height_px))
    fill = shape.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = RGBColor(*style["panel_fill"][0])
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = RGBColor(*style["panel_fill"][1])
    fill.gradient_stops[1].position = 1.0
    shape.fill.transparency = style.get("panel_transparency", transparency)
    shape.line.color.rgb = accent if style.get("panel_border") == "accent" else RGBColor(255, 255, 255)
    shape.line.width = Pt(1.2)
    add_glow_to_shape(shape, accent, size=style["panel_glow"])
    add_outer_shadow(shape, color_rgb="%02X%02X%02X" % (accent[0], accent[1], accent[2]),
                     blur_rad=style["panel_shadow_blur"], dist=12700, direction=5400000, alpha_pct=style["panel_shadow_alpha"])
    return shape


CANVAS_PRESETS = {
    "widescreen": {
        "width": 1920,
        "height": 1080,
        "slide_w": Inches(13.333333),
        "slide_h": Inches(7.5),
    },
    "xhs-vertical": {
        "width": 1080,
        "height": 1440,
        "slide_w": Inches(7.5),
        "slide_h": Inches(10),
    },
    "lecture-vertical": {
        "width": 1080,
        "height": 1920,
        "slide_w": Inches(7.5),
        "slide_h": Inches(13.333333),
    },
}

import platform
import os
import glob

def _find_font(candidates: list[str]) -> str | None:
    """Return the first existing font path from a list of candidates."""
    for p in candidates:
        if os.path.exists(p):
            return p
    return None

def _fallback_cjk_font() -> str:
    """Find any CJK-capable font on the system via fc-list."""
    try:
        result = subprocess.run(
            ["fc-list", ":lang=zh", "file"],
            capture_output=True, text=True, timeout=5,
        )
        for line in result.stdout.strip().splitlines():
            path = line.split(":")[0].strip()
            if path and os.path.exists(path):
                return path
    except Exception:
        pass
    raise RuntimeError(
        "No CJK font found. Install Noto Sans CJK or another CJK font, "
        "or add your font path to _FONT_CANDIDATES_BLACK."
    )

# --- Font resolution with cross-platform fallbacks ---
_FONT_CANDIDATES_BLACK = [
    "/usr/share/fonts/opentype/noto/NotoSansCJK-Black.ttc",   # Linux
    "/System/Library/Fonts/STHeiti Medium.ttc",                 # macOS
    "/Library/Fonts/Arial Unicode.ttf",                          # macOS fallback
    "/System/Library/Fonts/Supplemental/Songti.ttc",             # macOS fallback
]

_FONT_CANDIDATES_MONO = [
    "/usr/share/fonts/truetype/dejavu/DejaVuSansMono.ttf",
    "/System/Library/Fonts/Menlo.ttc",
    "/System/Library/Fonts/Monaco.ttf",
]

FONT_PATH_BLACK = _find_font(_FONT_CANDIDATES_BLACK) or _fallback_cjk_font()
FONT_PATH_MONO = _find_font(_FONT_CANDIDATES_MONO) or _FONT_CANDIDATES_MONO[-1]

COLORS = {
    "WHITE": RGBColor(255, 255, 255),
    "MUTED": RGBColor(188, 194, 210),
    "SOFT": RGBColor(120, 132, 154),
    "CARD": RGBColor(10, 10, 10),
    "CARD_2": RGBColor(5, 5, 8),
    "CYAN": RGBColor(0, 255, 255),
    "BLUE": RGBColor(59, 130, 246),
    "ORANGE": RGBColor(249, 115, 22),
    "YELLOW": RGBColor(251, 191, 36),
    "PINK": RGBColor(236, 72, 153),
    "RED": RGBColor(255, 51, 102),
    "PURPLE": RGBColor(139, 92, 246),
    "LIME": RGBColor(16, 185, 129),
    "TEAL": RGBColor(20, 184, 166),
    "AMBER": RGBColor(245, 158, 11),
    "CORAL": RGBColor(251, 113, 88),
    "PEACH": RGBColor(253, 186, 116),
    "ROSE": RGBColor(244, 114, 182),
    "GOLD": RGBColor(250, 204, 21),
}

STYLE_PRESETS = {
    "classic-cyberpunk": {
        "description": "Default dark neon cyberpunk visual system.",
        "page_label": "POSTER MODE / CYBERPUNK PPT",
        "accent_cycle": [
            ("RED", "YELLOW", "CYAN"),
            ("CYAN", "PURPLE", "PINK"),
            ("YELLOW", "TEAL", "ORANGE"),
            ("BLUE", "PINK", "LIME"),
        ],
        "background": (0, 0, 0, 255),
        "grid_color": (255, 255, 255),
        "frame_color": (255, 255, 255),
        "glow_alpha": (14, 11, 12),
        "glow_blur": 35,
        "grid_min": 60,
        "grid_divisor": 24,
        "grid_alpha": (12, 10),
        "grid_blur": 2,
        "frame_alpha": 25,
        "panel_fill": ((10, 10, 18), (18, 18, 32)),
        "panel_transparency": 0.30,
        "panel_border": "white",
        "panel_glow": 40000,
        "panel_shadow_blur": 50000,
        "panel_shadow_alpha": 25000,
        "chip_fill": ((10, 10, 18), (20, 20, 30)),
        "chip_transparency": 0.25,
        "chip_glow": 30000,
        "chip_shadow_alpha": 20000,
        "title_glow": 54000,
        "title_scale": 0.46,
        "title_min_pt": 26,
        "subtitle_size": 18,
        "line_glow": 18000,
        "tag_accent": "CYAN",
        "course_frame": False,
    },
    "warm-cyber": {
        "description": "Editorial warm cyberpunk style with amber glow and dense architecture diagrams.",
        "page_label": "WARM CYBER / SOFT NEON",
        "accent_cycle": [
            ("AMBER", "PEACH", "GOLD"),
            ("PEACH", "CORAL", "AMBER"),
            ("GOLD", "AMBER", "TEAL"),
            ("CORAL", "PEACH", "GOLD"),
        ],
        "background": (13, 8, 8, 255),
        "grid_color": (239, 164, 102),
        "frame_color": (255, 210, 150),
        "glow_alpha": (9, 7, 5),
        "glow_blur": 64,
        "grid_min": 96,
        "grid_divisor": 34,
        "grid_alpha": (9, 6),
        "grid_blur": 1,
        "frame_alpha": 38,
        "panel_fill": ((28, 16, 14), (42, 22, 16)),
        "panel_transparency": 0.30,
        "panel_border": "accent",
        "panel_glow": 4200,
        "panel_shadow_blur": 8000,
        "panel_shadow_alpha": 4200,
        "chip_fill": ((32, 18, 14), (50, 25, 16)),
        "chip_transparency": 0.16,
        "chip_glow": 4200,
        "chip_shadow_alpha": 8000,
        "title_glow": 9000,
        "title_scale": 0.36,
        "title_min_pt": 22,
        "subtitle_size": 16,
        "line_glow": 6000,
        "tag_accent": "AMBER",
        "warm_overlay": True,
        "background_assets": [
            "assets/backgrounds/warm-cyber/warm-cyber-bg-01.png",
            "assets/backgrounds/warm-cyber/warm-cyber-bg-02.png",
            "assets/backgrounds/warm-cyber/warm-cyber-bg-03.png",
            "assets/backgrounds/warm-cyber/warm-cyber-bg-04.png",
            "assets/backgrounds/warm-cyber/warm-cyber-bg-05.png",
            "assets/backgrounds/warm-cyber/warm-cyber-bg-06.png",
            "assets/backgrounds/warm-cyber/warm-cyber-bg-07.png",
            "assets/backgrounds/warm-cyber/warm-cyber-bg-08.png",
            "assets/backgrounds/warm-cyber/warm-cyber-bg-09.png",
            "assets/backgrounds/warm-cyber/warm-cyber-bg-10.png",
        ],
        "diagram_title_scale": 0.28,
        "diagram_title_min_pt": 20,
        "diagram_subtitle_size": 14,
    },
}

_ACTIVE_STYLE_NAME = "classic-cyberpunk"

# Safe area margins (px) — tag at top, page number at bottom
SLIDE_SAFE = {
    "widescreen": {"max_y": 980, "max_x": 1860, "top_y": 110},
    "xhs-vertical": {"max_y": 1380, "max_x": 1020, "top_y": 110},
    "lecture-vertical": {"max_y": 1860, "max_x": 1020, "top_y": 150},
}


def px(value: float):
    return Inches(value / 144)


def color(name: str) -> RGBColor:
    try:
        return COLORS[name.upper()]
    except KeyError as exc:
        raise ValueError(f"Unsupported color: {name}") from exc


def display_accent(accent_name: str, idx: int = 0) -> str:
    """Keep warm-cyber decks elegant by limiting saturated accent variety."""
    if _ACTIVE_STYLE_NAME != "warm-cyber":
        return accent_name
    refined = {
        "PINK": "PEACH",
        "ROSE": "CORAL",
        "YELLOW": "GOLD",
        "ORANGE": "CORAL",
        "RED": "CORAL",
        "PURPLE": "PEACH",
        "BLUE": "TEAL",
        "CYAN": "TEAL",
        "LIME": "GOLD",
    }
    base = refined.get((accent_name or "AMBER").upper(), accent_name)
    restrained_cycle = ["AMBER", "PEACH", "GOLD", "CORAL", "AMBER", "TEAL"]
    if base in {"TEAL", "CORAL"} and idx % 4 not in {0, 3}:
        return restrained_cycle[idx % len(restrained_cycle)]
    return base


def get_style_preset(style_name: str | None) -> dict:
    key = (style_name or "classic-cyberpunk").strip().lower()
    try:
        return STYLE_PRESETS[key]
    except KeyError as exc:
        valid = ", ".join(sorted(STYLE_PRESETS))
        raise ValueError(f"Unsupported style: {style_name}. Valid styles: {valid}") from exc


def set_active_style(style_name: str | None) -> None:
    global _ACTIVE_STYLE_NAME
    key = (style_name or "classic-cyberpunk").strip().lower()
    get_style_preset(key)
    _ACTIVE_STYLE_NAME = key


def active_style() -> dict:
    return get_style_preset(_ACTIVE_STYLE_NAME)


def style_palette(style: dict, idx: int) -> list[RGBColor]:
    cycle = style["accent_cycle"]
    names = cycle[idx % len(cycle)]
    return [color(name) for name in names]


def to_rgb(color_value: RGBColor) -> tuple[int, int, int]:
    return (color_value[0], color_value[1], color_value[2])


def pil_font(size: int) -> ImageFont.FreeTypeFont:
    return ImageFont.truetype(FONT_PATH_BLACK, size)


def _clamp(val: int, lo: int, hi: int) -> int:
    return max(lo, min(val, hi))


def _box_height_for_pt(pt_size: int) -> int:
    """Minimum textbox height (px) to fit a given pt font size without overflow."""
    return max(50, int(pt_size * 2.4))


def _line_advance_for_pt(pt_size: int) -> int:
    """Y advance (px) after rendering one title line at the given pt size."""
    return _box_height_for_pt(pt_size) + 8


def _resolve_font_path(font_name: str) -> str:
    if "Mono" in font_name or "mono" in font_name:
        return FONT_PATH_MONO
    return FONT_PATH_BLACK


def is_cjk(ch: str) -> bool:
    """Return True if ch is a CJK character (for word-wrapping purposes)."""
    cp = ord(ch)
    if 0x4E00 <= cp <= 0x9FFF: return True   # CJK Unified Ideographs
    if 0x3400 <= cp <= 0x4DBF: return True   # CJK Extension A
    if 0xF900 <= cp <= 0xFAFF: return True   # CJK Compatibility Ideographs
    if 0x3040 <= cp <= 0x309F: return True   # Hiragana
    if 0x30A0 <= cp <= 0x30FF: return True   # Katakana
    if 0xAC00 <= cp <= 0xD7AF: return True   # Hangul Syllables
    if 0x3000 <= cp <= 0x303F: return True   # CJK Symbols and Punctuation
    if 0xFF00 <= cp <= 0xFFEF: return True   # Fullwidth Forms
    if 0x20000 <= cp <= 0x2A6DF: return True # CJK Extension B
    if 0x2A700 <= cp <= 0x2B73F: return True # CJK Extension C
    if 0x2B740 <= cp <= 0x2B81F: return True # CJK Extension D
    if 0x2B820 <= cp <= 0x2CEAF: return True # CJK Extension E
    if 0x2CEB0 <= cp <= 0x2EBEF: return True # CJK Extension F
    if 0x2EBF0 <= cp <= 0x2EE5D: return True # CJK Extension I
    if 0x30000 <= cp <= 0x3134F: return True # CJK Extension G
    if 0x31350 <= cp <= 0x323AF: return True # CJK Extension H
    if 0x2F800 <= cp <= 0x2FA1F: return True # CJK Compat Ideographs Supplement
    return False


def measure_text(text: str, font_path: str, font_size_pt: int, max_width_px: int) -> dict:
    """Measure text with Pillow getbbox() and simulate word wrapping.

    Returns dict with lines, num_lines, total_height_px, max_width_px.
    Handles CJK characters (each char is a word boundary) and Latin words.
    """
    font = ImageFont.truetype(font_path, font_size_pt)

    # Split into tokens: CJK chars as individual tokens, Latin words grouped
    tokens: list[str] = []
    for ch in text:
        if is_cjk(ch):
            tokens.append(ch)
        elif ch in (' ', '\t'):
            if tokens and tokens[-1] != ' ':
                tokens.append(' ')
        else:
            if tokens and tokens[-1] not in (' ', '') and is_cjk(tokens[-1][-1]):
                tokens.append(ch)
            elif tokens and tokens[-1] not in (' ', ''):
                tokens[-1] += ch
            else:
                tokens.append(ch)

    lines: list[str] = []
    current_line = ""

    for token in tokens:
        if token == ' ':
            current_line += ' '
            continue
        test_line = current_line + token
        bbox = font.getbbox(test_line)
        text_width = bbox[2] - bbox[0]
        if text_width > max_width_px and current_line.strip():
            lines.append(current_line.strip())
            current_line = token
        else:
            current_line = test_line

    if current_line.strip():
        lines.append(current_line.strip())

    ascent, descent = font.getmetrics()
    line_height = int((ascent + descent) * 1.15)
    total_height = len(lines) * line_height
    max_line_width = max(
        font.getbbox(line)[2] - font.getbbox(line)[0]
        for line in lines
    ) if lines else 0

    return {
        "lines": lines,
        "num_lines": len(lines),
        "total_height_px": total_height,
        "max_width_px": max_line_width,
        "line_height_px": line_height,
    }


def fit_text_to_box(
    text: str,
    font_path: str,
    max_width_px: int,
    max_height_px: int,
    max_pt: int = 18,
    min_pt: int = 8,
) -> int:
    """Return the largest font size (pt) that fits text in the box."""
    for pt in range(max_pt, min_pt - 1, -1):
        metrics = measure_text(text, font_path, pt, max_width_px)
        if metrics["total_height_px"] <= max_height_px:
            return pt
    return min_pt


def add_warm_cyber_background_details(
    img: Image.Image,
    idx: int,
    width: int,
    height: int,
    palette: list[RGBColor],
) -> Image.Image:
    """Add restrained warm cyber geometry so pages differ without visual noise."""
    a1, a2, a3 = [to_rgb(item) for item in palette]
    layer = Image.new("RGBA", (width, height), (0, 0, 0, 0))
    draw = ImageDraw.Draw(layer, "RGBA")
    variant = idx % 4
    left_safe = int(width * 0.06)
    right_safe = int(width * 0.94)

    def circuit(start_x: int, start_y: int, segments: list[tuple[int, int]], col: tuple[int, int, int], alpha: int = 48):
        x, y = start_x, start_y
        for dx, dy in segments:
            nx, ny = x + dx, y + dy
            draw.line((x, y, nx, ny), fill=col + (alpha,), width=2)
            draw.ellipse((nx - 4, ny - 4, nx + 4, ny + 4), fill=col + (alpha + 18,))
            x, y = nx, ny

    if variant == 0:
        panel = (int(width * 0.58), 72, width - 88, height - 128)
        draw.rectangle(panel, fill=a1 + (18,))
        for inset, alpha in [(0, 42), (28, 24), (56, 16)]:
            draw.rectangle((panel[0] + inset, panel[1] + inset, panel[2] - inset, panel[3] - inset), outline=a2 + (alpha,), width=2)
        draw.line((int(width * 0.60), 96, width - 120, 96), fill=a2 + (58,), width=2)
        draw.line((int(width * 0.60), height - 160, width - 120, height - 160), fill=a3 + (38,), width=2)
        circuit(width - 520, 170, [(110, 0), (0, 54), (160, 0), (0, 70)], a3, 38)
    elif variant == 1:
        draw.polygon(
            [
                (0, int(height * 0.62)),
                (width, int(height * 0.45)),
                (width, int(height * 0.58)),
                (0, int(height * 0.75)),
            ],
            fill=a1 + (22,),
        )
        draw.polygon(
            [
                (0, int(height * 0.76)),
                (width, int(height * 0.59)),
                (width, int(height * 0.64)),
                (0, int(height * 0.82)),
            ],
            fill=a2 + (12,),
        )
        for x in range(180, width - 160, 260):
            draw.line((x, int(height * 0.52), x + 104, int(height * 0.50)), fill=a2 + (46,), width=2)
            draw.line((x + 20, int(height * 0.70), x + 94, int(height * 0.69)), fill=a3 + (30,), width=1)
        circuit(left_safe + 80, int(height * 0.78), [(170, 0), (0, -42), (130, 0)], a3, 36)
    elif variant == 2:
        draw.rectangle((88, 118, width - 92, 184), fill=a2 + (18,))
        draw.rectangle((88, height - 218, width - 92, height - 158), fill=a1 + (15,))
        draw.line((108, 146, width - 112, 146), fill=a3 + (36,), width=1)
        draw.line((108, height - 188, width - 112, height - 188), fill=a2 + (34,), width=1)
        for y in range(238, height - 260, 86):
            draw.line((width - 470, y, width - 128, y), fill=a3 + (34,), width=1)
            draw.rectangle((width - 150, y - 5, width - 128, y + 5), fill=a3 + (36,))
        circuit(120, 260, [(130, 0), (0, 60), (190, 0), (0, 52), (120, 0)], a1, 34)
    else:
        draw.rectangle((86, 96, 440, 105), fill=a1 + (54,))
        draw.rectangle((86, 116, 300, 121), fill=a2 + (38,))
        draw.rectangle((right_safe - 420, height - 218, right_safe, height - 204), fill=a2 + (44,))
        draw.rectangle((right_safe - 260, height - 186, right_safe, height - 179), fill=a3 + (34,))
        for offset in range(0, 240, 48):
            draw.arc((right_safe - 380 - offset, 150 + offset, right_safe - 140 + offset, 390 + offset), 205, 292, fill=a1 + (20,), width=2)
        circuit(right_safe - 520, 190, [(100, 0), (0, 48), (130, 0), (0, 48), (96, 0)], a2, 36)

    return Image.alpha_composite(img, layer)


def get_canvas(spec: dict) -> dict:
    canvas_name = spec.get("canvas", "widescreen")
    try:
        return CANVAS_PRESETS[canvas_name]
    except KeyError as exc:
        raise ValueError(f"Unsupported canvas: {canvas_name}") from exc


# ---------------------------------------------------------------------------
# Background generation (unchanged)
# ---------------------------------------------------------------------------

def build_background(idx: int, slide_spec: dict, asset_dir: Path, width: int, height: int) -> Path:
    canvas_name = slide_spec.get("_canvas_name", "widescreen")
    if canvas_name == "lecture-vertical":
        return build_lecture_background(idx, slide_spec, asset_dir, width, height)

    return build_poster_background(idx, slide_spec, asset_dir, width, height)


def build_asset_background(idx: int, style: dict, asset_dir: Path, width: int, height: int) -> Path | None:
    background_assets = style.get("background_assets") or []
    if not background_assets:
        return None

    src_rel = background_assets[idx % len(background_assets)]
    src = Path(__file__).resolve().parents[1] / src_rel
    if not src.exists():
        raise FileNotFoundError(f"Background asset not found: {src_rel}")

    asset_dir.mkdir(parents=True, exist_ok=True)
    img = Image.open(src).convert("RGB")
    img = ImageOps.fit(img, (width, height), method=Image.Resampling.LANCZOS, centering=(0.5, 0.5))
    output = asset_dir / f"poster_bg_{idx + 1:02d}.jpg"
    img.save(output, quality=92, optimize=True, progressive=True)
    return output


def build_poster_background(idx: int, slide_spec: dict, asset_dir: Path, width: int, height: int) -> Path:
    asset_dir.mkdir(parents=True, exist_ok=True)
    style = get_style_preset(slide_spec.get("_style"))
    asset_bg = build_asset_background(idx, style, asset_dir, width, height)
    if asset_bg:
        return asset_bg

    palette = style_palette(style, idx)
    a1, a2, a3 = to_rgb(palette[0]), to_rgb(palette[1]), to_rgb(palette[2])

    img = Image.new("RGBA", (width, height), style["background"])

    glow = Image.new("RGBA", (width, height), (0, 0, 0, 0))
    gdraw = ImageDraw.Draw(glow, "RGBA")
    cx, cy = width // 2, height // 2
    g1, g2, g3 = style["glow_alpha"]
    gdraw.ellipse((cx - int(width * 0.22), cy - int(height * 0.28), cx + int(width * 0.22), cy + int(height * 0.28)), fill=a1 + (g1,))
    gdraw.ellipse((int(width * 0.72), int(height * 0.08), width + 60, int(height * 0.42)), fill=a2 + (g2,))
    gdraw.ellipse((-60, int(height * 0.62), int(width * 0.28), height + 40), fill=a3 + (g3,))
    glow = glow.filter(ImageFilter.GaussianBlur(radius=style["glow_blur"]))
    img = Image.alpha_composite(img, glow)
    draw = ImageDraw.Draw(img, "RGBA")

    if style.get("warm_overlay"):
        img = add_warm_cyber_background_details(img, idx, width, height, palette)
        draw = ImageDraw.Draw(img, "RGBA")

    grid_layer = Image.new("RGBA", (width, height), (0, 0, 0, 0))
    gdraw_grid = ImageDraw.Draw(grid_layer, "RGBA")
    grid_step = max(style["grid_min"], width // style["grid_divisor"])
    grid_alpha_x, grid_alpha_y = style["grid_alpha"]
    grid_color = style["grid_color"]
    for x in range(0, width, grid_step):
        gdraw_grid.line((x, 0, x, height), fill=grid_color + (grid_alpha_x,), width=1)
    for y in range(0, height, grid_step):
        gdraw_grid.line((0, y, width, y), fill=grid_color + (grid_alpha_y,), width=1)
    grid_layer = grid_layer.filter(ImageFilter.GaussianBlur(radius=style["grid_blur"]))
    img = Image.alpha_composite(img, grid_layer)
    draw = ImageDraw.Draw(img, "RGBA")

    ghost = slide_spec.get("ghost", "")
    if ghost:
        ghost_layer = Image.new("RGBA", (width, height), (0, 0, 0, 0))
        ghost_draw = ImageDraw.Draw(ghost_layer)
        ghost_draw.text((width - 100, int(height * 0.18)), ghost, font=pil_font(max(140, min(width, height) // 5)), fill=a3 + (18,), anchor="ra")
        ghost_layer = ghost_layer.filter(ImageFilter.GaussianBlur(radius=3))
        img = Image.alpha_composite(img, ghost_layer)
        draw = ImageDraw.Draw(img, "RGBA")

    draw.rounded_rectangle((24, 24, width - 24, height - 24), radius=10, outline=style["frame_color"] + (style["frame_alpha"],), width=1)

    output = asset_dir / f"poster_bg_{idx + 1:02d}.jpg"
    img.convert("RGB").save(output, quality=90, optimize=True, progressive=True)
    return output


def add_lecture_scanlines(image: Image.Image) -> Image.Image:
    width, height = image.size
    overlay = Image.new("RGBA", image.size, (0, 0, 0, 0))
    draw = ImageDraw.Draw(overlay, "RGBA")
    for y in range(0, height, 6):
        alpha = 8 if y % 12 == 0 else 4
        draw.line((0, y, width, y), fill=(255, 255, 255, alpha), width=1)
    return Image.alpha_composite(image, overlay)


def draw_layered_glow(draw: ImageDraw.ImageDraw, center: tuple[int, int], radii: list[int], color_value: tuple[int, int, int], alphas: list[int]) -> None:
    x, y = center
    for radius, alpha in zip(radii, alphas):
        draw.ellipse((x - radius, y - radius, x + radius, y + radius), fill=color_value + (alpha,))


def add_lecture_orb(image: Image.Image, palette: list[tuple[int, int, int]]) -> Image.Image:
    width, height = image.size
    layer = Image.new("RGBA", image.size, (0, 0, 0, 0))
    draw = ImageDraw.Draw(layer, "RGBA")
    cx = width // 2
    cy = height - 210
    r = 88

    draw.ellipse((cx - 150, cy - 42, cx + 150, cy + 78), fill=(255, 175, 90, 30))
    draw.ellipse((cx - r, cy - r, cx + r, cy + r), fill=(230, 232, 238, 220))
    draw.ellipse((cx - r + 12, cy - r + 16, cx + r - 12, cy + r - 10), fill=(82, 86, 102, 170))
    draw.ellipse((cx - 46, cy - 58, cx + 18, cy - 8), fill=(255, 255, 255, 148))
    for idx, color_value in enumerate(palette):
        rr = 96 + idx * 52
        draw.ellipse((cx - rr, cy - rr, cx + rr, cy + rr), outline=color_value + (36,), width=2)
    return Image.alpha_composite(image, layer)


def build_lecture_background(idx: int, slide_spec: dict, asset_dir: Path, width: int, height: int) -> Path:
    asset_dir.mkdir(parents=True, exist_ok=True)
    style = get_style_preset(slide_spec.get("_style"))
    palette = style_palette(style, idx)
    a1, a2, a3 = to_rgb(palette[0]), to_rgb(palette[1]), to_rgb(palette[2])

    img = Image.new("RGBA", (width, height), style["background"])

    glow = Image.new("RGBA", (width, height), (0, 0, 0, 0))
    gdraw = ImageDraw.Draw(glow, "RGBA")
    draw_layered_glow(gdraw, (width // 2, 320), [180, 130, 80], a1, [14, 10, 7])
    draw_layered_glow(gdraw, (140, 960), [160, 110, 70], a2, [12, 9, 6])
    draw_layered_glow(gdraw, (880, 1240), [150, 100, 60], a3, [12, 8, 5])
    draw_layered_glow(gdraw, (width // 2, 1580), [110, 70], a1, [8, 5])
    glow = glow.filter(ImageFilter.GaussianBlur(radius=30))
    img = Image.alpha_composite(img, glow)

    grid_layer = Image.new("RGBA", (width, height), (0, 0, 0, 0))
    gdraw_grid = ImageDraw.Draw(grid_layer, "RGBA")
    grid_step = max(54, width // 20)
    for x in range(0, width, grid_step):
        gdraw_grid.line((x, 0, x, height), fill=style["grid_color"] + (10,), width=1)
    for y in range(0, height, grid_step):
        gdraw_grid.line((0, y, width, y), fill=style["grid_color"] + (8,), width=1)
    grid_layer = grid_layer.filter(ImageFilter.GaussianBlur(radius=2))
    img = Image.alpha_composite(img, grid_layer)

    img = add_lecture_scanlines(img)
    img = add_lecture_orb(img, [a1, a2, a3])

    draw = ImageDraw.Draw(img, "RGBA")
    draw.rounded_rectangle((20, 20, width - 20, height - 20), radius=8, outline=style["frame_color"] + (style["frame_alpha"],), width=1)

    output = asset_dir / f"lecture_bg_{idx + 1:02d}.png"
    img.save(output)
    return output


# ---------------------------------------------------------------------------
# Shared shape helpers
# ---------------------------------------------------------------------------

def add_textbox(slide, left, top, width, height, paragraphs, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, auto_fit=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = valign
    frame.margin_left = Pt(4)
    frame.margin_right = Pt(4)
    frame.margin_top = Pt(2)
    frame.margin_bottom = Pt(2)

    # Auto-fit: measure total text height and shrink if needed
    if auto_fit and paragraphs:
        box_width_px = int(width / 12700)
        box_height_px = int(height / 12700)
        if box_width_px > 20 and box_height_px > 10:
            full_text = "\n".join(p["text"] for p in paragraphs)
            font_name = paragraphs[0].get("font", "Noto Sans CJK SC")
            font_path = _resolve_font_path(font_name)
            max_pt = paragraphs[0].get("size", 18)
            inner_width = box_width_px - 10
            best_pt = fit_text_to_box(full_text, font_path, inner_width, box_height_px - 6, max_pt=max_pt, min_pt=8)
            if best_pt < max_pt:
                paragraphs = [{**p, "size": best_pt} for p in paragraphs]

    for idx, spec in enumerate(paragraphs):
        paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        paragraph.alignment = align
        paragraph.space_after = Pt(spec.get("space_after", 0))
        paragraph.line_spacing = spec.get("line_spacing", 1.0)
        run = paragraph.add_run()
        run.text = spec["text"]
        font = run.font
        font.name = spec.get("font", "Noto Sans CJK SC")
        font.size = Pt(spec.get("size", 18))
        font.bold = spec.get("bold", True)
        text_color = spec.get("color", COLORS["WHITE"])
        font.color.rgb = text_color
        glow_size = spec.get("glow", 0)
        if glow_size > 0:
            add_glow_to_run(run, text_color, size=glow_size)
    return box


def add_tag(slide, text, canvas_name="widescreen"):
    style = active_style()
    if canvas_name == "lecture-vertical":
        add_textbox(
            slide,
            px(240),
            px(110),
            px(600),
            px(28),
            [{"text": text, "size": 12, "bold": False, "color": COLORS["MUTED"]}],
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
        return

    if _ACTIVE_STYLE_NAME == "warm-cyber" and canvas_name == "widescreen":
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px(82), px(20), px(390), px(30))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*style["panel_fill"][0])
        shape.fill.transparency = 0.42
        tag_accent = color(style["tag_accent"])
        shape.line.color.rgb = tag_accent
        shape.line.width = Pt(0.6)
        add_textbox(slide, px(108), px(25), px(338), px(18), [{"text": text, "size": 10, "bold": False, "color": COLORS["MUTED"]}], auto_fit=True)
        return

    width_px = 430 if canvas_name == "widescreen" else 300
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px(82), px(60), px(width_px), px(40))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["CARD"]
    shape.fill.transparency = 0.30
    tag_accent = color(style["tag_accent"])
    shape.line.color.rgb = tag_accent
    shape.line.width = Pt(0.8)
    add_glow_to_shape(shape, tag_accent, size=style["chip_glow"])
    add_textbox(slide, px(108), px(69), px(width_px - 40), px(24), [{"text": text, "size": 12, "bold": False, "color": COLORS["MUTED"]}])


def add_page_no(slide, num, canvas_name="widescreen"):
    style = active_style()
    if canvas_name == "widescreen":
        add_textbox(slide, px(1760), px(995), px(80), px(32), [{"text": f"{num:02d}", "size": 15, "bold": True, "color": COLORS["SOFT"]}], align=PP_ALIGN.RIGHT)
        add_textbox(slide, px(112), px(995), px(420), px(24), [{"text": style["page_label"], "size": 11, "bold": False, "color": COLORS["SOFT"]}])
    elif canvas_name == "lecture-vertical":
        add_textbox(slide, px(930), px(1818), px(60), px(24), [{"text": f"{num:02d}", "size": 12, "bold": False, "color": COLORS["MUTED"]}], align=PP_ALIGN.RIGHT)
    else:
        add_textbox(slide, px(900), px(1350), px(80), px(32), [{"text": f"{num:02d}", "size": 15, "bold": True, "color": COLORS["SOFT"]}], align=PP_ALIGN.RIGHT)
        add_textbox(slide, px(90), px(1350), px(320), px(24), [{"text": "XHS / CYBERPUNK COVER", "size": 11, "bold": False, "color": COLORS["SOFT"]}])


# ---------------------------------------------------------------------------
# Title blocks — fixed sizing to prevent text overflow
# ---------------------------------------------------------------------------

def add_title_block(slide, title_lines, subtitle, left_px=118, top_px=168, width_px=980):
    """Render title lines + subtitle with measured text. Returns Y position after the entire block."""
    style = active_style()
    y = top_px
    for item in title_lines:
        pixel_size = int(item["size"])
        pt_size = max(style["title_min_pt"], int(pixel_size * style["title_scale"]))
        text_color = color(item["color"])
        # Measure actual text to determine box height
        metrics = measure_text(item["text"], FONT_PATH_BLACK, pt_size, width_px - 10)
        box_h = max(_box_height_for_pt(pt_size), metrics["total_height_px"] + 12)
        add_textbox(
            slide,
            px(left_px),
            px(y),
            px(width_px),
            px(box_h),
            [{"text": item["text"], "size": pt_size, "color": text_color, "glow": style["title_glow"]}],
        )
        y += box_h + 10
    if subtitle:
        sub_text = " ".join(subtitle)
        sub_size = style["subtitle_size"]
        sub_metrics = measure_text(sub_text, FONT_PATH_BLACK, sub_size, min(width_px, 860) - 10)
        sub_h = max(72, sub_metrics["total_height_px"] + 16)
        add_textbox(
            slide,
            px(left_px + 4),
            px(y + 10),
            px(min(width_px, 860)),
            px(sub_h),
            [{"text": sub_text, "size": sub_size, "bold": False, "color": COLORS["WHITE"], "line_spacing": 1.05}],
        )
        y += sub_h + 14
    return y


def add_title_block_vertical(slide, title_lines, subtitle, left_px=88, top_px=176, width_px=900):
    y = top_px
    for item in title_lines:
        pixel_size = int(item["size"])
        pt_size = max(24, int(pixel_size * 0.38))
        text_color = color(item["color"])
        metrics = measure_text(item["text"], FONT_PATH_BLACK, pt_size, width_px - 10)
        box_h = max(_box_height_for_pt(pt_size), metrics["total_height_px"] + 12)
        add_textbox(
            slide,
            px(left_px),
            px(y),
            px(width_px),
            px(box_h),
            [{"text": item["text"], "size": pt_size, "color": text_color, "glow": 48000}],
        )
        y += box_h + 10
    if subtitle:
        sub_text = " ".join(subtitle)
        sub_metrics = measure_text(sub_text, FONT_PATH_BLACK, 15, width_px - 30)
        sub_h = max(60, sub_metrics["total_height_px"] + 14)
        add_textbox(
            slide,
            px(left_px + 4),
            px(y + 8),
            px(width_px - 20),
            px(sub_h),
            [{"text": sub_text, "size": 15, "bold": False, "color": COLORS["WHITE"], "line_spacing": 1.02}],
        )
        y += sub_h + 10
    return y


def add_title_block_lecture(slide, title_lines, subtitle, top_px=260, width_px=820):
    y = top_px
    center_left = (1080 - width_px) // 2
    for item in title_lines:
        pixel_size = int(item["size"])
        pt_size = max(24, int(pixel_size * 0.36))
        text_color = color(item["color"])
        metrics = measure_text(item["text"], FONT_PATH_BLACK, pt_size, width_px - 10)
        box_h = max(_box_height_for_pt(pt_size), metrics["total_height_px"] + 12)
        add_textbox(
            slide,
            px(center_left),
            px(y),
            px(width_px),
            px(box_h),
            [{"text": item["text"], "size": pt_size, "color": text_color, "glow": 46000}],
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
        y += box_h + 10
    if subtitle:
        sub_text = " ".join(subtitle)
        sub_metrics = measure_text(sub_text, FONT_PATH_BLACK, 15, width_px - 30)
        sub_h = max(60, sub_metrics["total_height_px"] + 14)
        add_textbox(
            slide,
            px(center_left + 10),
            px(y + 14),
            px(width_px - 20),
            px(sub_h),
            [{"text": sub_text, "size": 15, "bold": False, "color": COLORS["WHITE"], "line_spacing": 1.05}],
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
        y += sub_h + 16
    return y


# ---------------------------------------------------------------------------
# Panel & chip helpers
# ---------------------------------------------------------------------------

def add_panel(slide, left_px, top_px, width_px, height_px, title, lines, accent_name, mono=False, title_size=18, body_size=16, canvas_name="widescreen"):
    accent = color(accent_name)
    safe = SLIDE_SAFE.get(canvas_name, SLIDE_SAFE["widescreen"])
    height_px = min(height_px, safe["max_y"] - top_px)
    if height_px < 60:
        height_px = 60
    add_gradient_panel(slide, left_px, top_px, width_px, height_px, accent_name)
    title_top = top_px + 18
    title_h = max(32, int(title_size * 2.2))
    add_accent_line(slide, left_px + 24, top_px + 14, min(width_px - 48, 72), accent_name, thickness=2)
    add_textbox(
        slide,
        px(left_px + 24),
        px(title_top),
        px(width_px - 48),
        px(title_h),
        [{"text": title, "size": title_size, "color": accent, "line_spacing": 1.0}],
        auto_fit=True,
    )
    body_font = "DejaVu Sans Mono" if mono else "Noto Sans CJK SC"
    effective_body_size = 14 if mono else body_size
    paragraphs = [{"text": line, "size": effective_body_size, "bold": False, "color": COLORS["WHITE"], "font": body_font, "space_after": 7} for line in lines]
    body_top = title_top + title_h + 10
    body_height = max(20, height_px - (body_top - top_px) - 18)
    add_textbox(slide, px(left_px + 24), px(body_top), px(width_px - 48), px(body_height), paragraphs, auto_fit=True)


def add_chip(slide, left_px, top_px, text, color_name):
    style = active_style()
    accent = color(color_name)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px(left_px), px(top_px), px(230), px(50))
    fill = shape.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = RGBColor(*style["chip_fill"][0])
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = RGBColor(*style["chip_fill"][1])
    fill.gradient_stops[1].position = 1.0
    shape.fill.transparency = style["chip_transparency"]
    shape.line.color.rgb = accent
    shape.line.width = Pt(1.0)
    add_glow_to_shape(shape, accent, size=style["chip_glow"])
    add_outer_shadow(shape, color_rgb="%02X%02X%02X" % (accent[0], accent[1], accent[2]),
                     blur_rad=22000, dist=12700, direction=5400000, alpha_pct=style["chip_shadow_alpha"])
    add_textbox(slide, px(left_px + 16), px(top_px + 12), px(198), px(22), [{"text": text, "size": 13, "color": accent}], align=PP_ALIGN.CENTER)


def add_timeline_label(slide, left_px, top_px, width_px, height_px, step: dict):
    accent = color(step["accent"])
    style = active_style()
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px(left_px), px(top_px), px(width_px), px(height_px))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*style["panel_fill"][0])
    shape.fill.transparency = min(0.28, style["panel_transparency"] + 0.08)
    shape.line.color.rgb = accent
    shape.line.width = Pt(0.8)
    add_glow_to_shape(shape, accent, size=max(4000, style["panel_glow"] // 2))
    add_textbox(
        slide,
        px(left_px + 16),
        px(top_px + 10),
        px(width_px - 32),
        px(height_px - 18),
        [{"text": step["label"], "size": 13, "color": accent, "line_spacing": 1.0}],
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        auto_fit=True,
    )


def add_info_box(slide, left_px, top_px, width_px, height_px, title, lines, accent_name, title_size=14, body_size=11, align=PP_ALIGN.LEFT):
    if _ACTIVE_STYLE_NAME == "warm-cyber":
        add_arch_box(
            slide,
            left_px,
            top_px,
            width_px,
            height_px,
            title,
            lines,
            display_accent(accent_name),
            title_size=title_size,
            body_size=body_size,
            align=align,
        )
        return
    add_panel(
        slide,
        left_px,
        top_px,
        width_px,
        height_px,
        title,
        lines,
        accent_name,
        title_size=title_size,
        body_size=body_size,
    )


def add_corner_ticks(slide, left_px, top_px, width_px, height_px, accent_name, tick=26):
    accent = color(accent_name)
    coords = [
        (left_px, top_px, tick, 1.2),
        (left_px, top_px, 1.2, tick),
        (left_px + width_px - tick, top_px, tick, 1.2),
        (left_px + width_px - 1.2, top_px, 1.2, tick),
        (left_px, top_px + height_px - 1.2, tick, 1.2),
        (left_px, top_px + height_px - tick, 1.2, tick),
        (left_px + width_px - tick, top_px + height_px - 1.2, tick, 1.2),
        (left_px + width_px - 1.2, top_px + height_px - tick, 1.2, tick),
    ]
    for x, y, w, h in coords:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(x), px(y), px(w), px(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = accent
        shape.fill.transparency = 0.12
        shape.line.fill.background()


def add_stroke_arrow_head(slide, x1_px, y1_px, x2_px, y2_px, accent_name, size_px=12, width=0.7):
    dx = x2_px - x1_px
    dy = y2_px - y1_px
    length = math.hypot(dx, dy)
    if length < 0.1:
        return
    ux = dx / length
    uy = dy / length
    pxn = -uy
    pyn = ux
    back_x = x2_px - ux * size_px
    back_y = y2_px - uy * size_px
    half = size_px * 0.46
    wing_1 = (back_x + pxn * half, back_y + pyn * half)
    wing_2 = (back_x - pxn * half, back_y - pyn * half)
    add_line_segment(slide, wing_1[0], wing_1[1], x2_px, y2_px, accent_name, width=width)
    add_line_segment(slide, wing_2[0], wing_2[1], x2_px, y2_px, accent_name, width=width)


def add_line_segment(slide, x1_px, y1_px, x2_px, y2_px, accent_name, width=0.7, dashed=False, end_arrow=False):
    accent = color(display_accent(accent_name))
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, px(x1_px), px(y1_px), px(x2_px), px(y2_px))
    line.line.color.rgb = accent
    line.line.width = Pt(width)
    if dashed:
        line.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    if end_arrow:
        add_stroke_arrow_head(slide, x1_px, y1_px, x2_px, y2_px, accent_name, size_px=12, width=width)
    return line


def add_flow_dot(slide, x_px, y_px, accent_name, size_px=6):
    accent = color(display_accent(accent_name))
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, px(x_px - size_px / 2), px(y_px - size_px / 2), px(size_px), px(size_px))
    dot.fill.solid()
    dot.fill.fore_color.rgb = accent
    dot.line.fill.background()
    add_glow_to_shape(dot, accent, size=2600)
    return dot


def add_arrow_head(slide, x_px, y_px, direction="right", accent_name="AMBER", size_px=13):
    accent = color(display_accent(accent_name))
    head = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_TRIANGLE,
        px(x_px - size_px / 2),
        px(y_px - size_px / 2),
        px(size_px),
        px(size_px),
    )
    rotations = {"right": 0, "down": 90, "left": 180, "up": 270}
    head.rotation = rotations.get(direction, 0)
    head.fill.solid()
    head.fill.fore_color.rgb = accent
    head.line.fill.background()
    return head


def add_flow_path(slide, points, accent_name="AMBER", width=0.72, arrow=True, dashed=False, dots=True):
    if len(points) < 2:
        return
    accent_name = display_accent(accent_name)
    segments = list(zip(points, points[1:]))
    for idx, (start, end) in enumerate(segments):
        add_line_segment(
            slide,
            start[0],
            start[1],
            end[0],
            end[1],
            accent_name,
            width=width,
            dashed=dashed,
            end_arrow=False,
        )
    if dots:
        for point in points[1:-1]:
            add_flow_dot(slide, point[0], point[1], accent_name, size_px=5)
    if arrow:
        start, end = segments[-1]
        add_stroke_arrow_head(slide, start[0], start[1], end[0], end[1], accent_name, size_px=11, width=width)


def add_smooth_flow_curve(slide, start, end, accent_name="AMBER", width=0.72, arrow=True, samples=18):
    accent_name = display_accent(accent_name)
    accent = color(accent_name)
    x1, y1 = start
    x2, y2 = end
    dx = x2 - x1
    control_span = dx * 0.42
    p1 = (x1 + control_span, y1)
    p2 = (x2 - control_span, y2)
    curve_points = []
    for step in range(samples + 1):
        t = step / samples
        inv = 1 - t
        x = inv**3 * x1 + 3 * inv**2 * t * p1[0] + 3 * inv * t**2 * p2[0] + t**3 * x2
        y = inv**3 * y1 + 3 * inv**2 * t * p1[1] + 3 * inv * t**2 * p2[1] + t**3 * y2
        curve_points.append((x, y))
    path = slide.shapes.build_freeform(px(x1), px(y1))
    path.add_line_segments([(px(x), px(y)) for x, y in curve_points[1:]], close=False)
    shape = path.convert_to_shape()
    shape.fill.background()
    shape.line.color.rgb = accent
    shape.line.width = Pt(width)
    if arrow:
        a, b = curve_points[-2], curve_points[-1]
        add_stroke_arrow_head(slide, a[0], a[1], b[0], b[1], accent_name, size_px=11, width=width)


def add_gateway_arrow(slide, left_px, top_px, width_px, height_px, label, accent_name="PEACH"):
    accent_name = display_accent(accent_name)
    y = top_px + height_px / 2
    add_flow_path(slide, [(left_px, y), (left_px + width_px - 12, y)], accent_name, width=1.0, arrow=True, dots=False)
    add_line_segment(slide, left_px + 12, top_px + 8, left_px + 12, top_px + height_px - 8, accent_name, width=0.55)
    add_textbox(
        slide,
        px(left_px + 10),
        px(top_px + height_px / 2 - 12),
        px(max(40, width_px - 24)),
        px(24),
        [{"text": label, "size": 11, "color": COLORS["WHITE"]}],
        align=PP_ALIGN.CENTER,
        auto_fit=True,
    )


def add_hub_intake_bus(slide, row_centers, source_x, bus_x, hub_x, hub_y, accent_name="PEACH"):
    """Bundle dense row outputs into a clean manifold before they enter the hub."""
    if not row_centers:
        return
    accent_name = display_accent(accent_name)
    top_y = min(row_centers) - 10
    bottom_y = max(row_centers) + 10
    add_line_segment(slide, bus_x, top_y, bus_x, bottom_y, accent_name, width=0.62, dashed=True)
    for y in row_centers:
        add_line_segment(slide, source_x, y, bus_x, y, accent_name, width=0.54)
        add_flow_dot(slide, bus_x, y, accent_name, size_px=4.4)
    for y in [hub_y - 56, hub_y, hub_y + 56]:
        add_flow_path(slide, [(bus_x, y), (hub_x - 92, y)], accent_name, width=0.78, arrow=True, dots=False)


def add_broken_frame(slide, left_px, top_px, width_px, height_px, accent_name, major=True):
    accent_name = display_accent(accent_name)
    inset = 18 if major else 12
    notch = 44 if major else 24
    lw = 0.72 if major else 0.55
    x1, y1 = left_px, top_px
    x2, y2 = left_px + width_px, top_px + height_px

    if major:
        segments = [
            ((x1 + inset, y1), (x1 + width_px * 0.42, y1)),
            ((x1 + width_px * 0.56, y1), (x2 - inset, y1)),
            ((x2, y1 + inset), (x2, y1 + height_px * 0.42)),
            ((x2, y1 + height_px * 0.58), (x2, y2 - inset)),
            ((x2 - inset, y2), (x1 + width_px * 0.58, y2)),
            ((x1 + width_px * 0.42, y2), (x1 + inset, y2)),
            ((x1, y2 - inset), (x1, y1 + height_px * 0.58)),
            ((x1, y1 + height_px * 0.42), (x1, y1 + inset)),
        ]
    else:
        segments = [
            ((x1 + inset, y1), (x2 - inset, y1)),
            ((x2, y1 + inset), (x2, y2 - inset)),
            ((x2 - inset, y2), (x1 + inset, y2)),
            ((x1, y2 - inset), (x1, y1 + inset)),
        ]
    for start, end in segments:
        add_line_segment(slide, start[0], start[1], end[0], end[1], accent_name, width=lw)

    cut_marks = [
        ((x1 + inset, y1), (x1 + inset + notch, y1)),
        ((x2 - inset - notch, y1), (x2 - inset, y1)),
        ((x1 + inset, y2), (x1 + inset + notch, y2)),
        ((x2 - inset - notch, y2), (x2 - inset, y2)),
    ]
    for start, end in cut_marks:
        add_line_segment(slide, start[0], start[1], end[0], end[1], accent_name, width=1.05 if major else 0.85)
    if major:
        add_flow_dot(slide, x2 - 2, y1 + 2, accent_name, size_px=4)


def add_arch_box(slide, left_px, top_px, width_px, height_px, title, lines, accent_name, title_size=14, body_size=11, align=PP_ALIGN.LEFT):
    accent_name = display_accent(accent_name)
    accent = color(accent_name)
    panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(left_px), px(top_px), px(width_px), px(height_px))
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(26, 15, 12)
    panel.fill.transparency = 0.30
    panel.line.fill.background()
    add_broken_frame(slide, left_px, top_px, width_px, height_px, accent_name, major=height_px > 120)
    compact = height_px <= 96
    accent_y = top_px + (12 if compact else 18)
    title_top = top_px + (20 if compact else 28)
    title_h = max(22, int(title_size * (1.7 if compact else 2.1)))
    body_top = title_top + title_h + (4 if compact else 8)
    add_accent_line(slide, left_px + 24, accent_y, min(width_px - 48, 64), accent_name, thickness=1.2)
    add_textbox(
        slide,
        px(left_px + 24),
        px(title_top),
        px(width_px - 48),
        px(title_h),
        [{"text": title, "size": title_size, "color": accent, "line_spacing": 1.0}],
        align=align,
        auto_fit=True,
    )
    paragraphs = [
        {"text": line, "size": body_size, "bold": False, "color": COLORS["WHITE"], "space_after": 5, "line_spacing": 1.0}
        for line in lines
    ]
    add_textbox(
        slide,
        px(left_px + 24),
        px(body_top),
        px(width_px - 48),
        px(max(20, height_px - (body_top - top_px) - 14)),
        paragraphs,
        align=align,
        auto_fit=True,
    )


def add_small_arrow(slide, left_px, top_px, width_px, color_name="AMBER"):
    color_name = display_accent(color_name)
    mid_y = top_px + 8
    if _ACTIVE_STYLE_NAME == "warm-cyber":
        add_flow_path(slide, [(left_px, mid_y), (left_px + width_px - 7, mid_y)], color_name, width=0.75, arrow=True, dots=False)
        add_flow_dot(slide, left_px, mid_y, color_name, size_px=3.5)
        return
    add_line_segment(slide, left_px, mid_y, left_px + width_px - 8, mid_y, color_name, width=0.95, end_arrow=True)


def add_hub(slide, cx_px, cy_px, radius_px, label, accent_name="AMBER"):
    accent_name = display_accent(accent_name)
    accent = color(accent_name)
    for idx, alpha in enumerate([0.72, 0.84, 0.92]):
        size = radius_px * (2 + idx * 0.72)
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, px(cx_px - size / 2), px(cy_px - size / 2), px(size), px(size))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(28, 16, 12)
        shape.fill.transparency = alpha
        shape.line.color.rgb = accent
        shape.line.width = Pt(0.7)
    core = slide.shapes.add_shape(MSO_SHAPE.OVAL, px(cx_px - radius_px), px(cy_px - radius_px), px(radius_px * 2), px(radius_px * 2))
    core.fill.solid()
    core.fill.fore_color.rgb = RGBColor(44, 24, 16)
    core.fill.transparency = 0.10
    core.line.color.rgb = accent
    core.line.width = Pt(1.4)
    add_glow_to_shape(core, accent, size=14000)
    add_textbox(
        slide,
        px(cx_px - radius_px + 18),
        px(cy_px - 34),
        px(radius_px * 2 - 36),
        px(68),
        [{"text": label, "size": 21, "color": accent, "line_spacing": 1.0}],
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        auto_fit=True,
    )


def add_diagram_title_block(slide, title_lines, subtitle, left_px=52, top_px=36, width_px=780):
    """Compact title block for dense warm-cyber architecture diagrams."""
    style = active_style()
    scale = style.get("diagram_title_scale", style["title_scale"])
    min_pt = style.get("diagram_title_min_pt", style["title_min_pt"])
    y = top_px
    for idx, item in enumerate(title_lines):
        max_pt = 44 if idx == 0 else 28
        pt_size = min(max_pt, max(min_pt, int(int(item["size"]) * scale)))
        text_color = color(item["color"])
        metrics = measure_text(item["text"], FONT_PATH_BLACK, pt_size, width_px - 10)
        box_h = max(_box_height_for_pt(pt_size), metrics["total_height_px"] + 8)
        add_textbox(
            slide,
            px(left_px),
            px(y),
            px(width_px),
            px(box_h),
            [{"text": item["text"], "size": pt_size, "color": text_color, "glow": style["title_glow"]}],
            auto_fit=True,
        )
        y += box_h + (4 if idx == 0 else 8)
    if subtitle:
        sub_text = " ".join(subtitle)
        sub_size = style.get("diagram_subtitle_size", style["subtitle_size"])
        sub_metrics = measure_text(sub_text, FONT_PATH_BLACK, sub_size, width_px - 24)
        sub_h = max(36, sub_metrics["total_height_px"] + 10)
        add_textbox(
            slide,
            px(left_px + 4),
            px(y + 2),
            px(width_px - 24),
            px(sub_h),
            [{"text": sub_text, "size": sub_size, "bold": False, "color": COLORS["WHITE"], "line_spacing": 1.0}],
            auto_fit=True,
        )
        y += sub_h + 8
    add_accent_line(slide, left_px + 2, y, min(width_px, 460), "PEACH", thickness=1.4)
    return y + 12


def add_connector_line(slide, x1_px, y1_px, x2_px, y2_px, color_name="AMBER", width=1.0):
    color_name = display_accent(color_name)
    if _ACTIVE_STYLE_NAME == "warm-cyber":
        add_flow_path(slide, [(x1_px, y1_px), (x2_px, y2_px)], color_name, width=max(0.55, width), arrow=False, dots=False)
        return None
    accent = color(color_name)
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        px(x1_px),
        px(y1_px),
        px(x2_px),
        px(y2_px),
    )
    line.line.color.rgb = accent
    line.line.width = Pt(width)
    return line


def add_number_badge(slide, left_px, top_px, number, accent_name="AMBER", size_px=46):
    accent_name = display_accent(accent_name, number)
    accent = color(accent_name)
    ring = slide.shapes.add_shape(MSO_SHAPE.OVAL, px(left_px), px(top_px), px(size_px), px(size_px))
    ring.fill.solid()
    ring.fill.fore_color.rgb = RGBColor(30, 17, 12)
    ring.fill.transparency = 0.18
    ring.line.color.rgb = accent
    ring.line.width = Pt(1.1)
    add_glow_to_shape(ring, accent, size=7000)
    add_textbox(
        slide,
        px(left_px + 3),
        px(top_px + 8),
        px(size_px - 6),
        px(size_px - 14),
        [{"text": str(number).zfill(2), "size": 15, "color": accent, "line_spacing": 1.0}],
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        auto_fit=True,
    )


def add_micro_visual(slide, left_px, top_px, width_px, height_px, accent_name="AMBER", variant=0):
    accent_name = display_accent(accent_name, variant)
    accent = color(accent_name)
    frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px(left_px), px(top_px), px(width_px), px(height_px))
    frame.fill.solid()
    frame.fill.fore_color.rgb = RGBColor(18, 12, 10)
    frame.fill.transparency = 0.48 if _ACTIVE_STYLE_NAME == "warm-cyber" else 0.34
    frame.line.color.rgb = accent
    frame.line.width = Pt(0.38 if _ACTIVE_STYLE_NAME == "warm-cyber" else 0.55)

    if variant % 3 == 0:
        for idx, h in enumerate([0.28, 0.54, 0.38, 0.72, 0.48]):
            bar_w = max(8, (width_px - 34) / 5)
            x = left_px + 12 + idx * (bar_w + 3)
            y = top_px + height_px - 10 - height_px * h
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(x), px(y), px(bar_w), px(height_px * h))
            bar.fill.solid()
            bar.fill.fore_color.rgb = accent
            bar.fill.transparency = 0.48 if _ACTIVE_STYLE_NAME == "warm-cyber" else 0.24
            bar.line.fill.background()
    elif variant % 3 == 1:
        points = [
            (left_px + 18, top_px + height_px - 18),
            (left_px + width_px * 0.34, top_px + 18),
            (left_px + width_px * 0.58, top_px + height_px * 0.62),
            (left_px + width_px - 18, top_px + height_px * 0.30),
        ]
        for a, b in zip(points, points[1:]):
            add_connector_line(slide, a[0], a[1], b[0], b[1], accent_name, width=0.75)
        for x, y in points:
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, px(x - 4), px(y - 4), px(8), px(8))
            dot.fill.solid()
            dot.fill.fore_color.rgb = accent
            dot.line.fill.background()
    else:
        cell = 10
        cols = max(4, int((width_px - 18) // cell))
        rows = max(3, int((height_px - 16) // cell))
        for r in range(rows):
            for c in range(cols):
                if (r + c + variant) % 3 == 0:
                    tile = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        px(left_px + 10 + c * cell),
                        px(top_px + 8 + r * cell),
                        px(cell - 3),
                        px(cell - 3),
                    )
                    tile.fill.solid()
                    tile.fill.fore_color.rgb = accent
                    tile.fill.transparency = 0.30
                    tile.line.fill.background()


def add_process_row(slide, left_px, top_px, width_px, height_px, card, idx):
    accent_name = display_accent(card.get("accent", "AMBER"), idx)
    accent = color(accent_name)
    panel_shape = MSO_SHAPE.RECTANGLE if _ACTIVE_STYLE_NAME == "warm-cyber" else MSO_SHAPE.ROUNDED_RECTANGLE
    panel = slide.shapes.add_shape(panel_shape, px(left_px), px(top_px), px(width_px), px(height_px))
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(32, 17, 13)
    panel.fill.transparency = 0.34 if _ACTIVE_STYLE_NAME == "warm-cyber" else 0.20
    if _ACTIVE_STYLE_NAME == "warm-cyber":
        panel.line.fill.background()
        add_broken_frame(slide, left_px, top_px, width_px, height_px, accent_name, major=False)
    else:
        panel.line.color.rgb = accent
        panel.line.width = Pt(0.9)
        add_glow_to_shape(panel, accent, size=6000)

    add_number_badge(slide, left_px + 16, top_px + 15, idx + 1, accent_name, size_px=44)
    title = card.get("title", "")
    title_text = re.sub(r"^\d+\s*", "", title)
    add_textbox(
        slide,
        px(left_px + 76),
        px(top_px + 14),
        px(230),
        px(24),
        [{"text": title_text, "size": 14, "color": accent, "line_spacing": 1.0}],
        auto_fit=True,
    )
    add_textbox(
        slide,
        px(left_px + 76),
        px(top_px + 42),
        px(232),
        px(height_px - 50),
        [{"text": " / ".join(card.get("lines", [])[:3]), "size": 9, "bold": False, "color": COLORS["WHITE"], "line_spacing": 1.05}],
        auto_fit=True,
    )
    visual_w = 108 if _ACTIVE_STYLE_NAME == "warm-cyber" else 118
    visual_h = height_px - 24
    visual_count = 2 if _ACTIVE_STYLE_NAME == "warm-cyber" else 3
    visual_x = left_px + width_px - (visual_count * visual_w + (visual_count - 1) * 22 + 24)
    for vidx in range(visual_count):
        x = visual_x + vidx * (visual_w + 22)
        add_micro_visual(slide, x, top_px + 12, visual_w, visual_h, accent_name, variant=idx + vidx)
        if vidx < visual_count - 1:
            add_small_arrow(slide, x + visual_w + 2, top_px + height_px / 2 - 7, 16, accent_name)
    add_small_arrow(slide, left_px + 316, top_px + height_px / 2 - 8, 52, accent_name)


def add_step_ribbon(slide, steps, left_px, top_px, box_w=92, box_h=42, gap=38):
    for idx, step in enumerate(steps[:7]):
        accent_name = step.get("accent", "AMBER")
        add_timeline_label(slide, left_px + idx * (box_w + gap), top_px, box_w, box_h, step)
        if idx < min(len(steps), 7) - 1:
            add_small_arrow(slide, left_px + idx * (box_w + gap) + box_w + 5, top_px + 14, gap - 6, accent_name)


def add_bottom_legend(slide, left_px=48, top_px=970, width_px=1760):
    style = active_style()
    accent = color("AMBER")
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px(left_px), px(top_px), px(width_px), px(54))
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(*style["panel_fill"][0])
    panel.fill.transparency = 0.42
    panel.line.color.rgb = accent
    panel.line.width = Pt(0.55)
    items = [
        ("数据流", "AMBER"), ("控制流", "PEACH"), ("资源流", "GOLD"),
        ("反馈流", "CORAL"), ("计算节点", "AMBER"), ("安全节点", "TEAL"),
    ]
    x = left_px + 54
    for label, accent_name in items:
        add_small_arrow(slide, x, top_px + 20, 70, accent_name)
        add_textbox(
            slide,
            px(x + 84),
            px(top_px + 16),
            px(128),
            px(22),
            [{"text": label, "size": 10, "bold": False, "color": COLORS["MUTED"]}],
            auto_fit=True,
        )
        x += 270


# ---------------------------------------------------------------------------
# Widescreen renderers — dynamic positioning, boundary-aware
# ---------------------------------------------------------------------------

def render_system_map(slide, spec):
    add_diagram_title_block(slide, spec["title"], spec.get("subtitle", []), left_px=48, top_px=58, width_px=720)
    cards = spec.get("cards", [])
    left_cards = cards[:3]
    center_cards = cards[3:9]
    right_cards = cards[9:13]

    steps = spec.get("steps") or [
        {"label": re.sub(r"^\d+\s*", "", card.get("title", ""))[:4], "accent": card.get("accent", "AMBER")}
        for card in center_cards[:6]
    ]
    add_step_ribbon(slide, steps, 780, 78, box_w=90, box_h=40, gap=42)

    rail_y = 184
    rail_h = 642
    add_info_box(slide, 38, rail_y, 250, rail_h, spec.get("left_title", "输入层"), [], "AMBER", title_size=18, body_size=10)
    for idx, card in enumerate(left_cards):
        add_info_box(slide, 70, rail_y + 78 + idx * 166, 188, 132, card["title"], card.get("lines", []), card["accent"], title_size=13, body_size=10)

    add_gateway_arrow(slide, 286, 448, 70, 50, "输入", "PEACH")

    row_x = 360
    row_w = 790
    row_h = 80
    row_centers = []
    for idx, card in enumerate(center_cards[:6]):
        y = 174 + idx * 88
        add_process_row(slide, row_x, y, row_w, row_h, card, idx)
        row_centers.append(y + 37)
        add_small_arrow(slide, row_x + row_w + 10, y + 30, 48, card.get("accent", "AMBER"))

    hub_x = 1326
    hub_y = 446
    add_hub(slide, hub_x, hub_y, 90, spec.get("hub", "协同\n核心"), "AMBER")
    add_hub_intake_bus(slide, row_centers, row_x + row_w + 58, hub_x - 110, hub_x, hub_y, "PEACH")
    add_flow_path(slide, [(hub_x, 142), (hub_x, hub_y - 116)], "PEACH", width=0.58, arrow=True, dashed=True, dots=False)
    add_flow_path(slide, [(hub_x, hub_y + 116), (hub_x, 805)], "PEACH", width=0.58, arrow=True, dashed=True, dots=False)

    add_gateway_arrow(slide, 1510, 448, 74, 50, "输出", "PEACH")

    add_info_box(slide, 1622, rail_y, 250, rail_h, spec.get("right_title", "输出层"), [], "AMBER", title_size=18, body_size=10)
    for idx, card in enumerate(right_cards[:4]):
        y = rail_y + 66 + idx * 142
        add_info_box(slide, 1652, y, 188, 112, card["title"], card.get("lines", []), card["accent"], title_size=12, body_size=9)
        add_smooth_flow_curve(slide, (hub_x + 104, y + 40), (1640, y + 56), card["accent"], width=0.68, arrow=True)

    rows = spec.get("rows", [])[:4]
    start_x = 54
    bottom_y = 840
    row_w = 420
    for idx, row in enumerate(rows):
        x = start_x + idx * 456
        add_info_box(slide, x, bottom_y, row_w, 86, row["title"], [row["body"]], row["accent"], title_size=12, body_size=9)
    add_bottom_legend(slide, top_px=946)


def render_hub_spoke(slide, spec):
    add_diagram_title_block(slide, spec["title"], spec.get("subtitle", []), left_px=70, top_px=64, width_px=760)
    nodes = spec.get("nodes", [])[:8]
    cx, cy = 960, 538
    add_hub(slide, cx, cy, 118, spec.get("hub", "军团\n核心"), "AMBER")
    positions = [
        (86, 252), (456, 216), (1210, 216), (1570, 286),
        (86, 724), (456, 808), (1210, 808), (1570, 714),
    ]
    for node, (x, y) in zip(nodes, positions):
        w = 300 if x < 900 else 292
        add_info_box(slide, x, y, w, 118, node["title"], [node["body"]], node["accent"], title_size=14, body_size=10)
        node_mid_x = x + w if x < cx else x
        node_mid_y = y + 59
        hub_edge_x = cx - 135 if x < cx else cx + 135
        bend_x = cx - 205 if x < cx else cx + 205
        hub_y = cy + (node_mid_y - cy) * 0.28
        add_flow_path(slide, [(node_mid_x, node_mid_y), (bend_x, node_mid_y), (bend_x, hub_y), (hub_edge_x, hub_y)], node["accent"], width=0.68, arrow=True, dots=True)
        add_micro_visual(slide, x + w - 94, y + 18, 64, 54, node["accent"], variant=len(node["title"]))
    rows = spec.get("rows", [])[:3]
    if rows:
        for idx, row in enumerate(rows):
            add_info_box(slide, 290 + idx * 450, 932, 390, 70, row["title"], [row["body"]], row["accent"], title_size=11, body_size=8)
    else:
        add_bottom_legend(slide, top_px=946)


def render_pipeline_board(slide, spec):
    add_diagram_title_block(slide, spec["title"], spec.get("subtitle", []), left_px=64, top_px=64, width_px=780)
    steps = spec.get("steps", [])[:7]
    add_step_ribbon(slide, steps, 765, 112, box_w=94, box_h=44, gap=42)

    cards = spec.get("cards", [])[:8]
    if _ACTIVE_STYLE_NAME == "warm-cyber":
        positions = [
            (72, 252, 560, 246),
            (672, 252, 380, 246),
            (1092, 252, 342, 246),
            (1474, 252, 374, 246),
            (72, 552, 382, 236),
            (494, 552, 382, 236),
            (916, 552, 458, 236),
            (1414, 552, 434, 236),
        ]
        for idx, (card, pos) in enumerate(zip(cards, positions)):
            x, y, w, h = pos
            add_info_box(slide, x, y, w, h, card["title"], card.get("lines", []), card["accent"], title_size=15 if idx == 0 else 13, body_size=12 if idx == 0 else 11)
            if idx in {0, 2, 5, 7}:
                add_micro_visual(slide, x + w - 132, y + 30, 96, 70, card["accent"], variant=idx)
            if idx < len(cards) - 1 and idx not in {3, 7}:
                add_small_arrow(slide, x + w + 8, y + h / 2 - 8, max(24, positions[idx + 1][0] - x - w - 12), card["accent"])
    else:
        left_x = 72
        top_grid = 254
        card_w = 414
        card_h = 228
        for idx, card in enumerate(cards):
            row = idx // 4
            col = idx % 4
            x = left_x + col * 456
            y = top_grid + row * 258
            add_info_box(slide, x, y, card_w, card_h, card["title"], card.get("lines", []), card["accent"], title_size=15, body_size=12)
            add_micro_visual(slide, x + card_w - 140, y + 30, 104, 76, card["accent"], variant=idx)
            if col < 3:
                add_small_arrow(slide, x + card_w + 8, y + card_h / 2 - 8, 32, card["accent"])
    rows = spec.get("rows", [])[:4]
    for idx, row in enumerate(rows):
        add_info_box(slide, 72 + idx * 456, 842, 414, 78, row["title"], [row["body"]], row["accent"], title_size=11, body_size=9)
    add_bottom_legend(slide, top_px=958)


def render_dense_grid(slide, spec):
    safe = SLIDE_SAFE["widescreen"]
    bottom = add_title_block(slide, spec["title"], spec.get("subtitle", []), left_px=92, top_px=118, width_px=1320)
    add_accent_line(slide, 96, bottom + 4, 360, spec["title"][0].get("color", "AMBER"), thickness=2)

    cards = spec.get("cards", [])[:8]
    if not cards:
        return

    if _ACTIVE_STYLE_NAME == "warm-cyber" and len(cards) <= 6:
        grid_top = _clamp(bottom + 34, 306, 360)
        feature = cards[0]
        add_arch_box(slide, 86, grid_top, 540, 570, feature["title"], feature.get("lines", []), display_accent(feature["accent"], 0), title_size=18, body_size=15)
        add_micro_visual(slide, 410, grid_top + 50, 160, 110, feature["accent"], variant=0)
        positions = [
            (680, grid_top, 540, 260),
            (1274, grid_top, 560, 260),
            (680, grid_top + 310, 350, 260),
            (1078, grid_top + 310, 350, 260),
            (1476, grid_top + 310, 358, 260),
        ]
        for idx, (card, pos) in enumerate(zip(cards[1:], positions), start=1):
            x, y, w, h = pos
            add_arch_box(slide, x, y, w, h, card["title"], card.get("lines", []), display_accent(card["accent"], idx), title_size=15, body_size=12)
            if idx in {1, 2, 5}:
                add_micro_visual(slide, x + w - 124, y + 34, 88, 62, card["accent"], variant=idx)
        return

    cols = 3 if len(cards) <= 6 else 4
    rows = 2
    margin_x = 86
    gap_x = 28
    gap_y = 24
    grid_top = _clamp(bottom + 34, 300, 390)
    grid_bottom = safe["max_y"] - 18
    card_w = int((1920 - margin_x * 2 - gap_x * (cols - 1)) / cols)
    card_h = int((grid_bottom - grid_top - gap_y * (rows - 1)) / rows)

    for idx, card in enumerate(cards[: cols * rows]):
        row = idx // cols
        col = idx % cols
        x = margin_x + col * (card_w + gap_x)
        y = grid_top + row * (card_h + gap_y)
        title_size = 17 if cols == 3 else 15
        body_size = 14 if cols == 3 else 12
        if _ACTIVE_STYLE_NAME == "warm-cyber":
            add_arch_box(
                slide,
                x,
                y,
                card_w,
                card_h,
                card["title"],
                card.get("lines", []),
                display_accent(card["accent"], idx),
                title_size=title_size,
                body_size=body_size,
            )
        else:
            add_panel(
                slide,
                x,
                y,
                card_w,
                card_h,
                card["title"],
                card.get("lines", []),
                card["accent"],
                title_size=title_size,
                body_size=body_size,
            )
        if _ACTIVE_STYLE_NAME == "warm-cyber" and card_w > 360:
            add_micro_visual(slide, x + card_w - 128, y + 28, 92, 64, card["accent"], variant=idx)


def render_cover(slide, spec):
    safe = SLIDE_SAFE["widescreen"]
    bottom = add_title_block(slide, spec["title"], spec.get("subtitle", []), left_px=118, top_px=160, width_px=980)
    # Decorative accent line below title
    add_accent_line(slide, 118, bottom + 6, 320, "CYAN", thickness=3)
    # Card on the right, vertically centered relative to title
    cards = spec.get("cards", [])
    if cards:
        title_mid = (160 + bottom) // 2
        card_y = _clamp(title_mid - 100, 200, safe["max_y"] - 250)
        card_h = min(240, safe["max_y"] - card_y - 10)
        add_panel(slide, 1280, card_y, 520, card_h, cards[0]["title"], cards[0].get("lines", []), cards[0]["accent"])
    # Chips below title
    chip_y = _clamp(bottom + 20, 500, safe["max_y"] - 60)
    for i, chip in enumerate(spec.get("chips", [])[:4]):
        add_chip(slide, 120 + i * 255, chip_y, chip["text"], chip["color"])
    ghost = spec.get("ghost", "")
    if ghost:
        ghost_y = _clamp(bottom - 60, 300, safe["max_y"] - 80)
        add_textbox(slide, px(1400), px(ghost_y), px(340), px(100), [{"text": ghost, "size": 36, "color": COLORS["WHITE"], "glow": 30000}], align=PP_ALIGN.CENTER)


def render_poster_cards(slide, spec):
    safe = SLIDE_SAFE["widescreen"]
    bottom = add_title_block(slide, spec["title"], spec.get("subtitle", []), width_px=940)
    add_accent_line(slide, 118, bottom + 6, 280, "PINK", thickness=2)
    cards = spec.get("cards", [])
    base_y = _clamp(bottom + 50, 440, 660)
    card_h = min(220, safe["max_y"] - base_y - 20)
    positions = [(118, base_y + 20), (690, base_y), (1262, base_y + 20)]
    for card, (x, y) in zip(cards[:3], positions):
        add_panel(slide, x, y, 520, card_h, card["title"], card.get("lines", []), card["accent"])


def render_flow(slide, spec):
    safe = SLIDE_SAFE["widescreen"]
    bottom = add_title_block(slide, spec["title"], spec.get("subtitle", []), width_px=940)
    nodes = spec.get("nodes", [])
    base_y = _clamp(bottom + 40, 460, 600)
    node_w = 330
    gap = 40
    total = len(nodes[:4]) * node_w + max(0, len(nodes[:4]) - 1) * gap
    start_x = max(96, (1920 - total) // 2)
    arrows = ["CYAN", "PINK", "YELLOW"]
    node_h = min(148, safe["max_y"] - base_y - 20)
    for i, (node, x_off) in enumerate(zip(nodes[:4], range(len(nodes[:4])))):
        x = start_x + x_off * (node_w + gap)
        add_panel(slide, x, base_y, node_w, node_h, node["title"], [node["body"]], node["accent"])
        if i < min(3, len(nodes[:4]) - 1):
            arrow = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, px(x + node_w + 6), px(base_y + node_h // 2 - 20), px(28), px(40))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = color(arrows[i])
            arrow.line.color.rgb = color(arrows[i])


def render_grid_four(slide, spec):
    safe = SLIDE_SAFE["widescreen"]
    bottom = add_title_block(slide, spec["title"], spec.get("subtitle", []), width_px=980)
    cards = spec.get("cards", [])
    base_y = _clamp(bottom + 40, 460, 620)
    card_w = min(830, (1920 - 118 - 118 - 40) // 2)
    card_h = min(180, (safe["max_y"] - base_y - 20) // 2 - 10)
    col2_x = 118 + card_w + 40
    positions = [(118, base_y), (col2_x, base_y), (118, base_y + card_h + 14), (col2_x, base_y + card_h + 14)]
    for card, (x, y) in zip(cards[:4], positions):
        add_panel(slide, x, y, card_w, card_h, card["title"], card.get("lines", []), card["accent"])


def render_split(slide, spec):
    safe = SLIDE_SAFE["widescreen"]
    bottom = add_title_block(slide, spec["title"], spec.get("subtitle", []), width_px=980)
    base_y = _clamp(bottom + 40, 460, 620)
    panel_w = min(760, (1920 - 118 - 80 - 40) // 2)
    panel_h = min(250, safe["max_y"] - base_y - 20)
    left = spec["left"]
    right = spec["right"]
    add_panel(slide, 118, base_y, panel_w, panel_h, left["title"], left.get("lines", []), left["accent"])
    add_panel(slide, 118 + panel_w + 40, base_y, panel_w, panel_h, right["title"], right.get("lines", []), right["accent"])


def render_code_mix(slide, spec):
    safe = SLIDE_SAFE["widescreen"]
    bottom = add_title_block(slide, spec["title"], spec.get("subtitle", []), width_px=960)
    base_y = _clamp(bottom + 30, 420, 560)
    code_h = min(300, safe["max_y"] - base_y - 20)
    add_panel(slide, 118, base_y, 740, code_h, "目录 / 命令", spec.get("code", []), "CYAN", mono=True)
    cards = spec.get("cards", [])
    card_x = 920
    card_w = min(700, 1920 - card_x - 50)
    card_h = min(110, (safe["max_y"] - base_y - 20) // max(1, len(cards[:3])) - 10)
    for idx, card in enumerate(cards[:3]):
        cy = base_y + idx * (card_h + 12)
        add_panel(slide, card_x, cy, card_w, card_h, card["title"], card.get("lines", []), card["accent"])


def render_timeline(slide, spec):
    safe = SLIDE_SAFE["widescreen"]
    bottom = add_title_block(slide, spec["title"], spec.get("subtitle", []), width_px=960)
    steps = spec.get("steps", [])
    if not steps:
        return
    style = active_style()
    line_y = _clamp(bottom + 108, 610, 760)
    # Horizontal line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(180), px(line_y), px(1500), px(3))
    line.fill.solid()
    line.fill.fore_color.rgb = color(style["tag_accent"])
    line.fill.transparency = 0.30
    line.line.fill.background()
    n = len(steps[:6])
    gap = 1440 // max(1, n - 1) if n > 1 else 0
    xs = [240 + i * gap for i in range(n)]
    dot_y = line_y - 22
    for si, (step, x) in enumerate(zip(steps[:6], xs)):
        accent = color(step["accent"])
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, px(x - 18), px(dot_y + 4), px(36), px(36))
        dot.fill.solid()
        dot.fill.fore_color.rgb = accent
        dot.line.color.rgb = accent
        add_textbox(slide, px(x - 21), px(dot_y + 11), px(42), px(18), [{"text": step["num"], "size": 10, "color": COLORS["CARD_2"]}], align=PP_ALIGN.CENTER)
        label_y = dot_y - 86 if si % 2 == 0 else line_y + 48
        label_y = _clamp(label_y, bottom + 28, safe["max_y"] - 82)
        add_timeline_label(slide, x - 120, label_y, 240, 68, step)


def render_wide_stack(slide, spec):
    safe = SLIDE_SAFE["widescreen"]
    bottom = add_title_block(slide, spec["title"], spec.get("subtitle", []), width_px=980)
    rows = spec.get("rows", [])
    base_y = _clamp(bottom + 30, 440, 580)
    row_h = min(90, (safe["max_y"] - base_y - 10) // max(1, len(rows[:4])) - 6)
    for i, row in enumerate(rows[:4]):
        add_panel(slide, 118, base_y + i * (row_h + 8), 1680, row_h, row["title"], [row["body"]], row["accent"])


def render_statement(slide, spec):
    safe = SLIDE_SAFE["widescreen"]
    bottom = add_title_block(slide, spec["title"], spec.get("subtitle", []), left_px=150, top_px=220, width_px=930)
    lines = spec.get("lines", [])
    base_y = _clamp(bottom + 50, 500, safe["max_y"] - 80)
    n = len(lines[:4])
    if n == 0:
        return
    item_w = min(300, (1600) // n)
    total_w = n * item_w
    start_x = max(120, (1920 - total_w) // 2)
    for i, item in enumerate(lines[:4]):
        add_textbox(slide, px(start_x + i * item_w), px(base_y), px(item_w - 10), px(70), [{"text": item["text"], "size": 32, "color": color(item["color"])}], align=PP_ALIGN.CENTER)


def render_ending(slide, spec):
    bottom = add_title_block(slide, spec["title"], spec.get("subtitle", []), left_px=150, top_px=250, width_px=980)
    add_accent_line(slide, 150, bottom + 6, 400, "CYAN", thickness=2)
    add_textbox(slide, px(360), px(860), px(1200), px(30), [{"text": spec.get("footer", ""), "size": 12, "bold": False, "color": COLORS["MUTED"]}], align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# XHS vertical renderers
# ---------------------------------------------------------------------------

def render_cover_vertical(slide, spec):
    safe = SLIDE_SAFE["xhs-vertical"]
    bottom = add_title_block_vertical(slide, spec["title"], spec.get("subtitle", []), left_px=88, top_px=180, width_px=820)
    chips = spec.get("chips", [])
    chip_start_y = _clamp(bottom + 40, 400, safe["max_y"] - 200)
    for idx, chip in enumerate(chips[:4]):
        row = idx // 2
        col = idx % 2
        add_chip(slide, 88 + col * 300, chip_start_y + row * 74, chip["text"], chip["color"])
    cards = spec.get("cards", [])
    if cards:
        card = cards[0]
        card_y = _clamp(chip_start_y + len(chips[:4]) * 74 + 30, bottom + 100, safe["max_y"] - 200)
        card_h = min(180, safe["max_y"] - card_y - 20)
        add_panel(slide, 88, card_y, 900, card_h, card["title"], card.get("lines", []), card["accent"])
    ghost = spec.get("ghost", "")
    if ghost:
        ghost_y = safe["max_y"] - 160
        add_textbox(slide, px(720), px(ghost_y), px(240), px(100), [{"text": ghost, "size": 34, "color": COLORS["WHITE"], "glow": 28000}], align=PP_ALIGN.CENTER)


def render_poster_cards_vertical(slide, spec):
    safe = SLIDE_SAFE["xhs-vertical"]
    bottom = add_title_block_vertical(slide, spec["title"], spec.get("subtitle", []), left_px=88, top_px=176, width_px=900)
    start_y = _clamp(bottom + 40, 400, 600)
    card_h = min(150, (safe["max_y"] - start_y - 10) // max(1, len(spec.get("cards", [])[:3])) - 10)
    for idx, card in enumerate(spec.get("cards", [])[:3]):
        add_panel(slide, 88, start_y + idx * (card_h + 14), 900, card_h, card["title"], card.get("lines", []), card["accent"])


def render_flow_vertical(slide, spec):
    safe = SLIDE_SAFE["xhs-vertical"]
    bottom = add_title_block_vertical(slide, spec["title"], spec.get("subtitle", []), left_px=88, top_px=176, width_px=900)
    start_y = _clamp(bottom + 40, 380, 580)
    nodes = spec.get("nodes", [])[:4]
    node_h = min(118, (safe["max_y"] - start_y - 10) // max(1, len(nodes)) - 30)
    for idx, node in enumerate(nodes):
        y = start_y + idx * (node_h + 40)
        add_panel(slide, 130, y, 820, node_h, node["title"], [node["body"]], node["accent"])
        if idx < len(nodes) - 1:
            add_textbox(
                slide,
                px(490),
                px(y + node_h + 4),
                px(100),
                px(30),
                [{"text": "▼", "size": 22, "color": color(nodes[idx + 1]["accent"])}],
                align=PP_ALIGN.CENTER,
            )


def render_grid_four_vertical(slide, spec):
    safe = SLIDE_SAFE["xhs-vertical"]
    bottom = add_title_block_vertical(slide, spec["title"], spec.get("subtitle", []), left_px=88, top_px=176, width_px=900)
    start_y = _clamp(bottom + 40, 400, 600)
    card_w = (900 - 20) // 2
    card_h = min(180, (safe["max_y"] - start_y - 10) // 2 - 10)
    positions = [(88, start_y), (88 + card_w + 20, start_y), (88, start_y + card_h + 14), (88 + card_w + 20, start_y + card_h + 14)]
    for card, (x, y) in zip(spec.get("cards", [])[:4], positions):
        add_panel(slide, x, y, card_w, card_h, card["title"], card.get("lines", []), card["accent"])


def render_split_vertical(slide, spec):
    safe = SLIDE_SAFE["xhs-vertical"]
    bottom = add_title_block_vertical(slide, spec["title"], spec.get("subtitle", []), left_px=88, top_px=176, width_px=900)
    start_y = _clamp(bottom + 40, 420, 650)
    panel_h = min(180, (safe["max_y"] - start_y - 10) // 2 - 10)
    add_panel(slide, 88, start_y, 900, panel_h, spec["left"]["title"], spec["left"].get("lines", []), spec["left"]["accent"])
    add_panel(slide, 88, start_y + panel_h + 14, 900, panel_h, spec["right"]["title"], spec["right"].get("lines", []), spec["right"]["accent"])


def render_code_mix_vertical(slide, spec):
    safe = SLIDE_SAFE["xhs-vertical"]
    bottom = add_title_block_vertical(slide, spec["title"], spec.get("subtitle", []), left_px=88, top_px=176, width_px=900)
    code_y = _clamp(bottom + 30, 380, 560)
    code_h = min(220, safe["max_y"] - code_y - 10)
    add_panel(slide, 88, code_y, 900, code_h, "目录 / 命令", spec.get("code", []), "CYAN", mono=True)
    cards = spec.get("cards", [])
    card_h = min(108, (safe["max_y"] - code_y - code_h - 10) // max(1, len(cards[:3])) - 8)
    for idx, card in enumerate(cards[:3]):
        cy = code_y + code_h + 14 + idx * (card_h + 10)
        add_panel(slide, 88, cy, 900, card_h, card["title"], card.get("lines", []), card["accent"])


def render_timeline_vertical(slide, spec):
    safe = SLIDE_SAFE["xhs-vertical"]
    bottom = add_title_block_vertical(slide, spec["title"], spec.get("subtitle", []), left_px=88, top_px=176, width_px=900)
    start_y = _clamp(bottom + 40, 400, 600)
    steps = spec.get("steps", [])[:5]
    step_h = min(96, (safe["max_y"] - start_y - 10) // max(1, len(steps)) - 10)
    for idx, step in enumerate(steps):
        label = f"{step['num']}  {step['label']}"
        add_panel(slide, 140, start_y + idx * (step_h + 12), 800, step_h, label, [], step["accent"])


def render_wide_stack_vertical(slide, spec):
    safe = SLIDE_SAFE["xhs-vertical"]
    bottom = add_title_block_vertical(slide, spec["title"], spec.get("subtitle", []), left_px=88, top_px=176, width_px=900)
    start_y = _clamp(bottom + 40, 400, 600)
    rows = spec.get("rows", [])[:4]
    row_h = min(112, (safe["max_y"] - start_y - 10) // max(1, len(rows)) - 8)
    for idx, row in enumerate(rows):
        add_panel(slide, 88, start_y + idx * (row_h + 10), 900, row_h, row["title"], [row["body"]], row["accent"])


def render_statement_vertical(slide, spec):
    safe = SLIDE_SAFE["xhs-vertical"]
    bottom = add_title_block_vertical(slide, spec["title"], spec.get("subtitle", []), left_px=88, top_px=220, width_px=900)
    start_y = _clamp(bottom + 50, 500, safe["max_y"] - 200)
    for idx, item in enumerate(spec.get("lines", [])[:4]):
        add_textbox(
            slide,
            px(120),
            px(start_y + idx * 60),
            px(840),
            px(48),
            [{"text": item["text"], "size": 26, "color": color(item["color"])}],
            align=PP_ALIGN.CENTER,
        )


def render_ending_vertical(slide, spec):
    add_title_block_vertical(slide, spec["title"], spec.get("subtitle", []), left_px=88, top_px=220, width_px=900)
    add_textbox(
        slide,
        px(120),
        px(1200),
        px(840),
        px(42),
        [{"text": spec.get("footer", ""), "size": 14, "bold": False, "color": COLORS["MUTED"]}],
        align=PP_ALIGN.CENTER,
    )


# ---------------------------------------------------------------------------
# Lecture vertical renderers
# ---------------------------------------------------------------------------

def render_cover_lecture(slide, spec):
    safe = SLIDE_SAFE["lecture-vertical"]
    bottom = add_title_block_lecture(slide, spec["title"], spec.get("subtitle", []), top_px=330, width_px=840)
    chips = spec.get("chips", [])
    chip_y = _clamp(bottom + 30, 600, safe["max_y"] - 200)
    for idx, chip in enumerate(chips[:4]):
        row = idx // 2
        col = idx % 2
        add_chip(slide, 152 + col * 392, chip_y + row * 82, chip["text"], chip["color"])
    cards = spec.get("cards", [])
    if cards:
        card_y = _clamp(chip_y + len(chips[:4]) * 82 + 30, bottom + 100, safe["max_y"] - 200)
        card_h = min(190, safe["max_y"] - card_y - 20)
        add_panel(slide, 120, card_y, 840, card_h, cards[0]["title"], cards[0].get("lines", []), cards[0]["accent"], body_size=15)


def render_poster_cards_lecture(slide, spec):
    safe = SLIDE_SAFE["lecture-vertical"]
    bottom = add_title_block_lecture(slide, spec["title"], spec.get("subtitle", []), top_px=260, width_px=860)
    start_y = _clamp(bottom + 50, 600, 900)
    cards = spec.get("cards", [])[:3]
    card_h = min(132, (safe["max_y"] - start_y - 10) // max(1, len(cards)) - 10)
    for idx, card in enumerate(cards):
        add_panel(slide, 118, start_y + idx * (card_h + 14), 844, card_h, card["title"], card.get("lines", []), card["accent"], body_size=15)


def render_grid_four_lecture(slide, spec):
    safe = SLIDE_SAFE["lecture-vertical"]
    bottom = add_title_block_lecture(slide, spec["title"], spec.get("subtitle", []), top_px=260, width_px=860)
    start_y = _clamp(bottom + 50, 580, 860)
    card_w = (844 - 20) // 2
    card_h = min(164, (safe["max_y"] - start_y - 10) // 2 - 10)
    positions = [(108, start_y), (108 + card_w + 20, start_y), (108, start_y + card_h + 14), (108 + card_w + 20, start_y + card_h + 14)]
    for card, (x, y) in zip(spec.get("cards", [])[:4], positions):
        add_panel(slide, x, y, card_w, card_h, card["title"], card.get("lines", []), card["accent"], body_size=14)


def render_split_lecture(slide, spec):
    safe = SLIDE_SAFE["lecture-vertical"]
    bottom = add_title_block_lecture(slide, spec["title"], spec.get("subtitle", []), top_px=260, width_px=860)
    start_y = _clamp(bottom + 50, 600, 900)
    panel_h = min(184, (safe["max_y"] - start_y - 10) // 2 - 10)
    add_panel(slide, 118, start_y, 844, panel_h, spec["left"]["title"], spec["left"].get("lines", []), spec["left"]["accent"], body_size=14)
    add_panel(slide, 118, start_y + panel_h + 14, 844, panel_h, spec["right"]["title"], spec["right"].get("lines", []), spec["right"]["accent"], body_size=14)


def render_code_mix_lecture(slide, spec):
    safe = SLIDE_SAFE["lecture-vertical"]
    bottom = add_title_block_lecture(slide, spec["title"], spec.get("subtitle", []), top_px=260, width_px=860)
    start_y = _clamp(bottom + 40, 560, 820)
    code_h = min(250, safe["max_y"] - start_y - 10)
    add_panel(slide, 118, start_y, 844, code_h, "目录 / 命令", spec.get("code", []), "CYAN", mono=True)
    cards = spec.get("cards", [])
    card_h = min(128, (safe["max_y"] - start_y - code_h - 10) // max(1, len(cards[:3])) - 10)
    for idx, card in enumerate(cards[:3]):
        cy = start_y + code_h + 14 + idx * (card_h + 10)
        add_panel(slide, 118, cy, 844, card_h, card["title"], card.get("lines", []), card["accent"], body_size=13)


def render_flow_lecture(slide, spec):
    safe = SLIDE_SAFE["lecture-vertical"]
    bottom = add_title_block_lecture(slide, spec["title"], spec.get("subtitle", []), top_px=260, width_px=860)
    start_y = _clamp(bottom + 50, 600, 900)
    nodes = spec.get("nodes", [])[:4]
    node_h = min(106, (safe["max_y"] - start_y - 10) // max(1, len(nodes)) - 30)
    for idx, node in enumerate(nodes):
        y = start_y + idx * (node_h + 36)
        add_panel(slide, 162, y, 756, node_h, node["title"], [node["body"]], node["accent"], body_size=14)
        if idx < len(nodes) - 1:
            add_textbox(slide, px(500), px(y + node_h + 4), px(80), px(28), [{"text": "▼", "size": 18, "color": color(nodes[idx + 1]["accent"])}], align=PP_ALIGN.CENTER)


def render_timeline_lecture(slide, spec):
    safe = SLIDE_SAFE["lecture-vertical"]
    bottom = add_title_block_lecture(slide, spec["title"], spec.get("subtitle", []), top_px=260, width_px=860)
    start_y = _clamp(bottom + 50, 600, 900)
    steps = spec.get("steps", [])[:5]
    step_h = min(92, (safe["max_y"] - start_y - 10) // max(1, len(steps)) - 10)
    for idx, step in enumerate(steps):
        add_panel(slide, 128, start_y + idx * (step_h + 12), 824, step_h, f"{step['num']}  {step['label']}", [], step["accent"], body_size=14)


def render_wide_stack_lecture(slide, spec):
    safe = SLIDE_SAFE["lecture-vertical"]
    bottom = add_title_block_lecture(slide, spec["title"], spec.get("subtitle", []), top_px=260, width_px=860)
    start_y = _clamp(bottom + 50, 600, 900)
    rows = spec.get("rows", [])[:4]
    row_h = min(112, (safe["max_y"] - start_y - 10) // max(1, len(rows)) - 8)
    for idx, row in enumerate(rows):
        add_panel(slide, 118, start_y + idx * (row_h + 10), 844, row_h, row["title"], [row["body"]], row["accent"], body_size=14)


def render_statement_lecture(slide, spec):
    safe = SLIDE_SAFE["lecture-vertical"]
    bottom = add_title_block_lecture(slide, spec["title"], spec.get("subtitle", []), top_px=280, width_px=860)
    line_y = _clamp(bottom + 60, 700, safe["max_y"] - 200)
    for idx, item in enumerate(spec.get("lines", [])[:4]):
        add_textbox(
            slide,
            px(140),
            px(line_y + idx * 60),
            px(800),
            px(42),
            [{"text": item["text"], "size": 26, "color": color(item["color"])}],
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )


def render_ending_lecture(slide, spec):
    add_title_block_lecture(slide, spec["title"], spec.get("subtitle", []), top_px=320, width_px=860)
    add_textbox(
        slide,
        px(160),
        px(1180),
        px(760),
        px(80),
        [{"text": spec.get("footer", ""), "size": 14, "bold": False, "color": COLORS["MUTED"], "line_spacing": 1.05}],
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )


# ---------------------------------------------------------------------------
# Renderer registries
# ---------------------------------------------------------------------------

RENDERERS = {
    "cover": render_cover,
    "dense_grid": render_dense_grid,
    "system_map": render_system_map,
    "hub_spoke": render_hub_spoke,
    "pipeline_board": render_pipeline_board,
    "poster_cards": render_poster_cards,
    "flow": render_flow,
    "grid_four": render_grid_four,
    "split": render_split,
    "code_mix": render_code_mix,
    "timeline": render_timeline,
    "wide_stack": render_wide_stack,
    "statement": render_statement,
    "ending": render_ending,
}

VERTICAL_RENDERERS = {
    "cover": render_cover_vertical,
    "poster_cards": render_poster_cards_vertical,
    "flow": render_flow_vertical,
    "grid_four": render_grid_four_vertical,
    "split": render_split_vertical,
    "code_mix": render_code_mix_vertical,
    "timeline": render_timeline_vertical,
    "wide_stack": render_wide_stack_vertical,
    "statement": render_statement_vertical,
    "ending": render_ending_vertical,
}

LECTURE_VERTICAL_RENDERERS = {
    "cover": render_cover_lecture,
    "poster_cards": render_poster_cards_lecture,
    "flow": render_flow_lecture,
    "grid_four": render_grid_four_lecture,
    "split": render_split_lecture,
    "code_mix": render_code_mix_lecture,
    "timeline": render_timeline_lecture,
    "wide_stack": render_wide_stack_lecture,
    "statement": render_statement_lecture,
    "ending": render_ending_lecture,
}


# ---------------------------------------------------------------------------
# Main generation pipeline
# ---------------------------------------------------------------------------

def make_presentation(spec: dict, output_path: Path, asset_dir: Path) -> None:
    canvas = get_canvas(spec)
    canvas_name = spec.get("canvas", "widescreen")
    deck_style = spec.get("style", "classic-cyberpunk")
    get_style_preset(deck_style)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = canvas["slide_w"]
    prs.slide_height = canvas["slide_h"]

    for idx, slide_spec in enumerate(spec["slides"]):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide_spec = dict(slide_spec)
        slide_spec["_canvas_name"] = canvas_name
        slide_spec["_style"] = slide_spec.get("style", deck_style)
        set_active_style(slide_spec["_style"])
        bg = build_background(idx, slide_spec, asset_dir, canvas["width"], canvas["height"])
        slide.shapes.add_picture(str(bg), 0, 0, width=canvas["slide_w"], height=canvas["slide_h"])
        add_tag(slide, slide_spec.get("tag", f"CUT {idx + 1:02d}"), canvas_name=canvas_name)
        add_page_no(slide, idx + 1, canvas_name=canvas_name)

        layout_name = slide_spec["layout"]
        if canvas_name == "xhs-vertical":
            registry = VERTICAL_RENDERERS
        elif canvas_name == "lecture-vertical":
            registry = LECTURE_VERTICAL_RENDERERS
        else:
            registry = RENDERERS
        if layout_name not in registry:
            valid = sorted(registry.keys())
            raise ValueError(
                f"Unknown layout '{layout_name}' (slide {idx + 1}). "
                f"Valid layouts for '{canvas_name}': {valid}"
            )
        registry[layout_name](slide, slide_spec)

    prs.save(str(output_path))


def load_spec(spec_path: Path) -> dict:
    return json.loads(spec_path.read_text(encoding="utf-8"))


def export_pdf(pptx_path: Path, pdf_output: Path) -> None:
    if shutil.which("libreoffice") is None:
        raise RuntimeError(
            "libreoffice is required for PDF export but not found. "
            "Install it with: sudo apt install libreoffice-impress"
        )
    pdf_output.parent.mkdir(parents=True, exist_ok=True)
    with tempfile.TemporaryDirectory(prefix="cyberpunk-pdf-") as tmpdir:
        tmpdir_path = Path(tmpdir)
        subprocess.run(
            [
                "libreoffice",
                "--headless",
                "--convert-to",
                "pdf",
                str(pptx_path),
                "--outdir",
                str(tmpdir_path),
            ],
            check=True,
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
        )
        generated_pdf = tmpdir_path / f"{pptx_path.stem}.pdf"
        if not generated_pdf.exists():
            raise FileNotFoundError(f"PDF export failed for {pptx_path}")
        shutil.copy2(generated_pdf, pdf_output)


def main() -> None:
    parser = argparse.ArgumentParser(description="Generate editable cyberpunk PPT from JSON spec.")
    parser.add_argument("--spec", required=True, help="Path to JSON spec file.")
    parser.add_argument("--output", help="Output PPTX path. Omit to auto-organize under ~/ai-gen-ppt/.")
    parser.add_argument("--assets-dir", help="Directory for generated background assets.")
    parser.add_argument("--pdf-output", help="Optional output PDF path.")
    args = parser.parse_args()

    spec_path = Path(args.spec)
    spec = load_spec(spec_path)

    if args.output:
        output_path = Path(args.output)
        asset_dir = Path(args.assets_dir) if args.assets_dir else Path("generated_cyberpunk_assets")
    else:
        deck_title = extract_deck_title(spec)
        safe_title = sanitize_dirname(deck_title)
        out_dir = resolve_output_dir(deck_title)
        output_path = out_dir / f"{safe_title}.pptx"
        asset_dir = out_dir / "assets"
        spec_path = out_dir / "spec.json"
        spec_path.write_text(json.dumps(spec, ensure_ascii=False, indent=2), encoding="utf-8")

    make_presentation(spec, output_path, asset_dir)

    if args.pdf_output:
        export_pdf(output_path, Path(args.pdf_output))

    print(f"Generated {output_path} with {len(spec['slides'])} slides")


if __name__ == "__main__":
    main()
