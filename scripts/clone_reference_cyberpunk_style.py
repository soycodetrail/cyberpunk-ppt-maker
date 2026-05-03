#!/usr/bin/env python3
from __future__ import annotations

import argparse
from pathlib import Path
import re
import tempfile

from pptx import Presentation

from export_cyberpunk_images import export_images
from generate_cyberpunk_ppt import (
    extract_deck_title,
    export_pdf,
    make_presentation,
    resolve_output_dir,
    sanitize_dirname,
)
from markdown_to_cyberpunk_spec import parse_markdown_outline, write_spec


def infer_canvas(reference_pptx: Path) -> str:
    prs = Presentation(str(reference_pptx))
    ratio = prs.slide_width / prs.slide_height
    if ratio < 0.65:
        return "lecture-vertical"
    if ratio < 1:
        return "xhs-vertical"
    return "widescreen"


def infer_tag_prefix(reference_pptx: Path) -> str | None:
    prs = Presentation(str(reference_pptx))
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if re.match(r".+\s\d{2}$", text):
                    return re.sub(r"\s\d{2}$", "", text)
    return None


def clone_from_reference(reference_pptx: Path, content_markdown: Path) -> dict:
    spec = parse_markdown_outline(content_markdown.read_text(encoding="utf-8"))
    spec["canvas"] = infer_canvas(reference_pptx)
    tag_prefix = infer_tag_prefix(reference_pptx)
    if tag_prefix:
        for idx, slide in enumerate(spec["slides"], 1):
            slide["tag"] = f"{tag_prefix} {idx:02d}"
    return spec


def main() -> None:
    parser = argparse.ArgumentParser(description="Use an existing PPT as a style reference entrypoint for cyberpunk outputs.")
    parser.add_argument("--reference-pptx", required=True, help="Reference PPTX path.")
    parser.add_argument("--content-markdown", required=True, help="Markdown content outline to restyle.")
    parser.add_argument("--output-spec", help="Output JSON spec path. Omit to auto-organize under ~/ai-gen-ppt/.")
    parser.add_argument("--pptx-output", help="Optional PPTX output path.")
    parser.add_argument("--pdf-output", help="Optional PDF output path.")
    parser.add_argument("--png-dir", help="Optional PNG output directory.")
    parser.add_argument("--assets-dir", help="Optional generated assets directory.")
    args = parser.parse_args()

    reference_pptx = Path(args.reference_pptx)
    content_markdown = Path(args.content_markdown)
    spec = clone_from_reference(reference_pptx, content_markdown)

    if args.output_spec:
        output_spec = Path(args.output_spec)
        assets_dir = Path(args.assets_dir) if args.assets_dir else output_spec.parent / "generated_cyberpunk_assets"
    else:
        deck_title = extract_deck_title(spec)
        safe_title = sanitize_dirname(deck_title)
        out_dir = resolve_output_dir(deck_title)
        output_spec = out_dir / "spec.json"
        assets_dir = out_dir / "assets"

    write_spec(spec, output_spec)

    if args.output_spec:
        pptx_output = Path(args.pptx_output) if args.pptx_output else None
        pdf_output = Path(args.pdf_output) if args.pdf_output else None
        png_dir = Path(args.png_dir) if args.png_dir else None
    else:
        deck_title = extract_deck_title(spec)
        safe_title = sanitize_dirname(deck_title)
        out_dir = output_spec.parent
        pptx_output = out_dir / f"{safe_title}.pptx"
        pdf_output = out_dir / f"{safe_title}.pdf" if args.pdf_output else None
        png_dir = out_dir / "png" if args.png_dir else None

    if pptx_output:
        make_presentation(spec, pptx_output, assets_dir)
        if pdf_output:
            export_pdf(pptx_output, pdf_output)
    elif pdf_output:
        with tempfile.TemporaryDirectory(prefix="cyberpunk-clone-") as tmpdir:
            temp_ppt = Path(tmpdir) / "clone.pptx"
            make_presentation(spec, temp_ppt, assets_dir)
            export_pdf(temp_ppt, pdf_output)

    if png_dir:
        export_images(output_spec, png_dir, assets_dir=assets_dir, pptx_path=pptx_output)

    print(f"Cloned style from {reference_pptx} into {output_spec}")


if __name__ == "__main__":
    main()
